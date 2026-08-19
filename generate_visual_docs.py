"""
グラフ・図表画像を埋め込んだドキュメントを生成するスクリプト。
Vision AI（画像読み取り）の評価ベンチマーク用。

生成するドキュメント:
  - 売上実績ダッシュボード_2024上半期.pdf       棒グラフ・円グラフをPNG画像として埋め込み
  - 組織図_テックリードジャパン_2024年4月.pptx  組織図をPNG画像として埋め込み
  - 業務フロー図_カスタマーサクセス対応.docx    フロー図をPNG画像として埋め込み

いずれも、重要な数値・名前・構造情報が画像の中にのみ存在し、
テキスト抽出（PyPDF / python-docx / python-pptx）では取得できない。
"""

import io
from pathlib import Path

DOCS = Path(__file__).parent / "docs"
DOCS.mkdir(exist_ok=True)

# ─── matplotlib セットアップ ──────────────────────────────────────────────────
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import matplotlib.patches as mpatches

# 日本語フォント検索
import matplotlib.font_manager as fm

_jp_candidates = [
    'Hiragino Sans', 'Hiragino Kaku Gothic Pro', 'Hiragino Kaku Gothic ProN',
    'IPAexGothic', 'IPAGothic', 'Noto Sans CJK JP', 'Yu Gothic', 'MS Gothic',
    'Arial Unicode MS',
]
_available = {f.name for f in fm.fontManager.ttflist}
_jp_font = next((f for f in _jp_candidates if f in _available), None)
if _jp_font:
    plt.rcParams['font.family'] = [_jp_font, 'DejaVu Sans']
    plt.rcParams['axes.unicode_minus'] = False
    print(f"  日本語フォント: {_jp_font}")
else:
    print("  警告: 日本語フォントが見つかりません。ラベルが文字化けする場合があります。")


# ─── PDF 1: 売上実績ダッシュボード ───────────────────────────────────────────
def make_dashboard_pdf():
    """棒グラフ・円グラフをPNG画像として埋め込んだ売上ダッシュボードPDF。
    各月の正確な売上高・比率はグラフ画像の中にのみ存在する。
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont

    pdfmetrics.registerFont(UnicodeCIDFont('HeiseiKakuGo-W5'))

    # ── 棒グラフ生成 ──────────────────────────────────────────────────────────
    months = ["4月", "5月", "6月", "7月", "8月", "9月"]
    sales  = [12.3, 13.8, 14.5, 15.2, 13.6, 16.8]   # 億円（グラフ内にのみ存在）

    fig, ax = plt.subplots(figsize=(9, 5.0))
    colors = ['#2E74B5'] * 6
    colors[5] = '#C00000'   # 最高月（9月）を赤
    bars = ax.bar(months, sales, color=colors, width=0.6, edgecolor='white', linewidth=0.5)
    # バーの上にデータラベルを表示（Vision AIが読み取れる数値）
    for bar, val in zip(bars, sales):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                f'{val}億円', ha='center', va='bottom', fontsize=11, fontweight='bold',
                color='#222')
    ax.set_ylabel("売上高（億円）", fontsize=12)
    ax.set_title("月別売上高推移（2024年度 上半期）", fontsize=14, fontweight='bold', pad=10)
    ax.set_ylim(0, 21)
    ax.yaxis.grid(True, linestyle='--', alpha=0.5, color='#aaa')
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    buf_bar = io.BytesIO()
    fig.savefig(buf_bar, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    bar_bytes = buf_bar.getvalue()
    plt.close(fig)

    # ── 円グラフ生成 ──────────────────────────────────────────────────────────
    seg_labels = ["ITソリューション", "AIサービス", "保守サポート", "新規事業"]
    seg_sizes  = [42, 28, 18, 12]    # %（グラフ内にのみ存在）
    seg_colors = ['#2E74B5', '#ED7D31', '#A5A5A5', '#FFC000']

    fig2, ax2 = plt.subplots(figsize=(7, 5))
    wedges, texts, autotexts = ax2.pie(
        seg_sizes, labels=seg_labels, colors=seg_colors,
        autopct='%1.0f%%', startangle=90, explode=(0, 0.05, 0, 0.08),
        textprops={'fontsize': 10},
    )
    for at in autotexts:
        at.set_fontsize(11)
        at.set_fontweight('bold')
    ax2.set_title("事業部別売上比率（2024年度 上半期）", fontsize=13, fontweight='bold', pad=10)

    buf_pie = io.BytesIO()
    fig2.savefig(buf_pie, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    pie_bytes = buf_pie.getvalue()
    plt.close(fig2)

    # ── PDF 組み立て ─────────────────────────────────────────────────────────
    out_path = DOCS / "売上実績ダッシュボード_2024上半期.pdf"
    c = canvas.Canvas(str(out_path), pagesize=A4)
    W, H = A4

    # ページ1: タイトル・概要（テキストのみ）
    c.setFont("HeiseiKakuGo-W5", 22)
    c.drawCentredString(W / 2, H - 2.8*cm, "売上実績ダッシュボード")
    c.setFont("HeiseiKakuGo-W5", 15)
    c.drawCentredString(W / 2, H - 4.2*cm, "2024年度 上半期（4月〜9月）")
    c.setFillColorRGB(0.4, 0.4, 0.4)
    c.setFont("HeiseiKakuGo-W5", 11)
    c.drawCentredString(W / 2, H - 5.3*cm, "経営管理部　作成　2024年10月15日")
    c.setFillColorRGB(0, 0, 0)
    c.setFont("HeiseiKakuGo-W5", 12)
    overview = [
        "【資料概要】",
        "本ダッシュボードは2024年度上半期（4〜9月）の売上実績を可視化した経営資料です。",
        "月別売上推移および事業部別売上比率を次ページ以降のグラフにて確認できます。",
        "",
        "【上半期KPIサマリ】",
        "・上半期売上高合計　　 ：86.2億円",
        "・前年同期比成長率　　 ：+14.3%",
        "・月別・事業部別の内訳 ：次ページ以降のグラフを参照",
        "",
        "【構成】",
        "・p.2　月別売上高推移（棒グラフ）",
        "・p.3　事業部別売上比率（円グラフ）",
    ]
    y = H - 7.0*cm
    for line in overview:
        c.drawString(2.5*cm, y, line)
        y -= 0.65*cm

    c.showPage()

    # ページ2: 棒グラフ画像
    c.setFont("HeiseiKakuGo-W5", 14)
    c.setFillColorRGB(0, 0, 0)
    c.drawString(1.5*cm, H - 1.8*cm, "グラフ①　月別売上高推移")
    iw = W - 3*cm
    ih = iw * 4.5 / 9
    c.drawImage(ImageReader(io.BytesIO(bar_bytes)), 1.5*cm, H - 2.2*cm - ih, iw, ih)
    c.showPage()

    # ページ3: 円グラフ画像
    c.setFont("HeiseiKakuGo-W5", 14)
    c.setFillColorRGB(0, 0, 0)
    c.drawString(1.5*cm, H - 1.8*cm, "グラフ②　事業部別売上比率")
    pw = W - 5*cm
    ph = pw * 5 / 7
    c.drawImage(ImageReader(io.BytesIO(pie_bytes)), 2.5*cm, H - 2.2*cm - ph, pw, ph)
    c.save()
    print("  ✓ 売上実績ダッシュボード_2024上半期.pdf")


# ─── PPTX 4: 組織図（PNG画像埋め込み） ──────────────────────────────────────
def make_org_chart_pptx():
    """組織図をmatplotlibで描画しPNG画像として埋め込んだPPTX。
    役職名・氏名・部署構造はすべて画像の中にのみ存在する。
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as PptxColor

    # ── 組織図PNG生成 ────────────────────────────────────────────────────────
    # 15×9 の広いキャンバス。Y字型コネクタで重なりゼロを保証する。
    # Lv2 ボックス幅 1.9、間隔 0.25 → 7箱分: 7×1.9 + 6×0.25 = 14.8 ≦ 15
    fig, ax = plt.subplots(figsize=(15, 9))
    ax.set_xlim(0, 15)
    ax.set_ylim(0, 9)
    ax.axis('off')
    fig.patch.set_facecolor('#F8F9FA')

    # ── Y レベル（ボックス中心） ───────────────────────────────────────────────
    Y0, Y1, Y2, Y3 = 8.3, 6.9, 5.3, 3.8
    BH = 0.6   # ボックス高さ（全階層共通）

    # ── X 中心座標 ────────────────────────────────────────────────────────────
    # Lv2 ボックスを左から均等配置し、グループごとに Lv1 の中心を決める
    # CFO配下 (2箱): x = 1.1, 3.3   →  CFO  x = 2.2
    # CTO配下 (3箱): x = 5.5, 7.5, 9.5   →  CTO  x = 7.5
    # 営業配下 (2箱): x = 11.7, 13.9  →  営業  x = 12.8
    XF1, XF2       = 1.1, 3.3
    XT1, XT2, XT3  = 5.5, 7.5, 9.5
    XS1, XS2       = 11.7, 13.9

    X_CFO  = (XF1 + XF2) / 2          # 2.2
    X_CTO  = (XT1 + XT2 + XT3) / 3   # 7.5
    X_SALE = (XS1 + XS2) / 2          # 12.8
    X_PRES = (X_CFO + X_CTO + X_SALE) / 3  # ≈ 7.43 → 7.5 に丸める
    X_PRES = 7.5

    # Lv3 (AI開発部直下) — XT1=5.5 の左右に均等配置
    XA1, XA2 = 4.5, 6.5

    # ── ボックス幅 ────────────────────────────────────────────────────────────
    W0 = 2.8   # 社長
    W1 = 2.4   # 執行役員
    W2 = 1.9   # 部 (Lv2)   ← 間隔 0.25 で重なりゼロ確認済み
    W3 = 1.7   # チーム (Lv3)

    # ── ヘルパー: ボックス描画 ────────────────────────────────────────────────
    def rbox(x, y, w, h, line1, line2='', bg='#2E74B5', fg='white', fs=8.5):
        rect = FancyBboxPatch((x - w/2, y - h/2), w, h,
                              boxstyle="round,pad=0.08",
                              facecolor=bg, edgecolor='#333', linewidth=1.2)
        ax.add_patch(rect)
        if line2:
            ax.text(x, y + 0.13, line1, ha='center', va='center',
                    color=fg, fontsize=fs - 1, style='italic')
            ax.text(x, y - 0.15, line2, ha='center', va='center',
                    color=fg, fontsize=fs + 0.5, fontweight='bold')
        else:
            ax.text(x, y, line1, ha='center', va='center',
                    color=fg, fontsize=fs, fontweight='bold',
                    multialignment='center')

    # ── ヘルパー: Y字型ツリーコネクタ ────────────────────────────────────────
    def tree_conn(px, py_bot, children_x, child_y_top, color='#555'):
        """親ボックス底辺 → 中継Y → 子ボックス上辺 を Y字で繋ぐ"""
        mid_y = (py_bot + child_y_top) / 2
        # 親から中継点へ
        ax.plot([px, px], [py_bot, mid_y], color=color, lw=1.4)
        # 中継点の水平線
        ax.plot([min(children_x), max(children_x)], [mid_y, mid_y],
                color=color, lw=1.4)
        # 各子への矢印
        for cx in children_x:
            ax.annotate('', xy=(cx, child_y_top), xytext=(cx, mid_y),
                        arrowprops=dict(arrowstyle='->', color=color, lw=1.4))

    # ── Lv0: 代表取締役社長 ───────────────────────────────────────────────────
    rbox(X_PRES, Y0, W0, BH, '代表取締役社長', '田中 一郎', bg='#1F3864', fs=9.5)

    # 社長 → 3役員
    tree_conn(X_PRES, Y0 - BH/2, [X_CFO, X_CTO, X_SALE], Y1 + BH/2)

    # ── Lv1: 執行役員 ────────────────────────────────────────────────────────
    rbox(X_CFO,  Y1, W1, BH, '専務取締役 CFO', '鈴木 二郎',   bg='#2E74B5', fs=8.5)
    rbox(X_CTO,  Y1, W1, BH, '常務取締役 CTO', '佐藤 三郎',   bg='#2E74B5', fs=8.5)
    rbox(X_SALE, Y1, W1, BH, '取締役 営業本部長', '高橋 四郎', bg='#2E74B5', fs=8.5)

    # CFO → 財務部・経理部
    tree_conn(X_CFO,  Y1 - BH/2, [XF1, XF2],       Y2 + BH/2)
    # CTO → AI開発部・インフラ部・セキュリティ部
    tree_conn(X_CTO,  Y1 - BH/2, [XT1, XT2, XT3],  Y2 + BH/2)
    # 営業 → フィールドSE・カスタマーS
    tree_conn(X_SALE, Y1 - BH/2, [XS1, XS2],        Y2 + BH/2)

    # ── Lv2: 部 ──────────────────────────────────────────────────────────────
    # CFO配下（緑）
    rbox(XF1, Y2, W2, BH, '財務部',        '部長: 木下 健',    bg='#70AD47', fs=8)
    rbox(XF2, Y2, W2, BH, '経理部',        '部長: 森田 純子',  bg='#70AD47', fs=8)
    # CTO配下（オレンジ）
    rbox(XT1, Y2, W2, BH, 'AI開発部',      '部長: 高田 健一',  bg='#ED7D31', fs=8)
    rbox(XT2, Y2, W2, BH, 'インフラ部',    '部長: 中村 誠',    bg='#ED7D31', fs=8)
    rbox(XT3, Y2, W2, BH, 'セキュリティ部', '部長: 山口 真理', bg='#ED7D31', fs=7.5)
    # 営業配下（紫）
    rbox(XS1, Y2, W2, BH, 'フィールドSE部', '部長: 河野 誠',   bg='#7030A0', fs=7.5)
    rbox(XS2, Y2, W2, BH, 'カスタマーS部',  '部長: 岡田 美咲', bg='#7030A0', fs=7.5)

    # AI開発部 → AIエンジニアG・MLOps G
    tree_conn(XT1, Y2 - BH/2, [XA1, XA2], Y3 + BH/2)

    # ── Lv3: チーム ──────────────────────────────────────────────────────────
    rbox(XA1, Y3, W3, BH, 'AIエンジニアG', '12名', bg='#FFC000', fg='#111', fs=7.5)
    rbox(XA2, Y3, W3, BH, 'MLOps G',       '8名',  bg='#FFC000', fg='#111', fs=7.5)

    # ── 凡例 ─────────────────────────────────────────────────────────────────
    legend_items = [
        mpatches.Patch(color='#1F3864', label='経営トップ'),
        mpatches.Patch(color='#2E74B5', label='執行役員'),
        mpatches.Patch(color='#70AD47', label='財務部門'),
        mpatches.Patch(color='#ED7D31', label='技術部門（CTO管轄）'),
        mpatches.Patch(color='#7030A0', label='営業部門'),
        mpatches.Patch(color='#FFC000', label='チーム'),
    ]
    ax.legend(handles=legend_items, loc='lower left', fontsize=8,
              framealpha=0.85, ncol=3, bbox_to_anchor=(0.0, 0.0))

    ax.set_title('テックリードジャパン 組織図（2024年4月1日改訂）',
                 fontsize=14, fontweight='bold', pad=10)

    buf_org = io.BytesIO()
    fig.savefig(buf_org, format='png', dpi=150, bbox_inches='tight', facecolor='#F8F9FA')
    org_bytes = buf_org.getvalue()
    plt.close(fig)

    # ── PPTX 組み立て ────────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # スライド1: タイトル（テキスト）
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "テックリードジャパン 組織図"
    s1.placeholders[1].text = "2024年4月1日改訂版\n人事部 作成"

    # スライド2: 組織図PNG
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s2.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(12.7), Inches(0.55))
    tf = tb.text_frame
    tf.text = "組織図 — 2024年4月1日現在"
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = PptxColor(0x1F, 0x38, 0x64)
    s2.shapes.add_picture(io.BytesIO(org_bytes),
                          Inches(0.15), Inches(0.65), Inches(13.0), Inches(6.7))

    # スライド3: テキスト補足（読み取れる情報）
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    tb3 = s3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(6.8))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p0 = tf3.paragraphs[0]
    p0.text = "組織概要"
    p0.runs[0].font.size = Pt(18)
    p0.runs[0].font.bold = True

    for text in [
        "代表取締役社長 田中一郎のもと、CFO・CTO・営業本部長の3名の執行役員が各部門を統括しています。",
        "技術部門（CTO管轄）、財務部門（CFO管轄）、営業部門（営業本部長管轄）の3ブロック体制です。",
        "各部門長・部署名・チーム人数の詳細は前スライドの組織図画像を参照してください。",
        "2024年4月1日付で組織再編を実施。AI開発部は同年2月設立の新設部門です。",
    ]:
        p = tf3.add_paragraph()
        p.text = text
        p.runs[0].font.size = Pt(13)
        p.space_after = Pt(8)

    prs.save(DOCS / "組織図_テックリードジャパン_2024年4月.pptx")
    print("  ✓ 組織図_テックリードジャパン_2024年4月.pptx")


# ─── DOCX 4: 業務フロー図（PNG画像埋め込み） ────────────────────────────────
def make_flowchart_docx():
    """フロー図をmatplotlibで描画しPNG画像として埋め込んだDOCX。
    対応時間・エスカレーション条件・判定ロジックはフロー図画像の中にのみ存在する。
    """
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.shared import RGBColor

    # ── フロー図PNG生成 ───────────────────────────────────────────────────────
    fig, ax = plt.subplots(figsize=(10, 10))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    C_START  = '#1F3864'
    C_PROC   = '#2E74B5'
    C_DEC    = '#ED7D31'
    C_CRIT   = '#C00000'
    C_HIGH   = '#E36C09'
    C_NORM   = '#4472C4'
    C_END    = '#70AD47'

    def rbox(x, y, w, h, text, color, fg='white', fs=8.5):
        rect = FancyBboxPatch((x - w/2, y - h/2), w, h,
                              boxstyle="round,pad=0.1",
                              facecolor=color, edgecolor='#222', linewidth=1.3)
        ax.add_patch(rect)
        ax.text(x, y, text, ha='center', va='center', color=fg,
                fontsize=fs, fontweight='bold', multialignment='center')

    def diamond(x, y, w, h, text, color, fg='white', fs=8.5):
        import numpy as np
        pts = [[x, y + h/2], [x + w/2, y], [x, y - h/2], [x - w/2, y]]
        poly = plt.Polygon(pts, facecolor=color, edgecolor='#222', linewidth=1.3)
        ax.add_patch(poly)
        ax.text(x, y, text, ha='center', va='center', color=fg,
                fontsize=fs, fontweight='bold', multialignment='center')

    def arrow(ax, x1, y1, x2, y2, label='', lx=0.15, ly=0):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color='#444', lw=1.4))
        if label:
            mx = (x1 + x2) / 2 + lx
            my = (y1 + y2) / 2 + ly
            ax.text(mx, my, label, fontsize=8, color='#444', va='center', fontweight='bold')

    # フローステップ
    rbox(5, 9.5, 4.5, 0.65,
         "顧客からの問い合わせ受付\n（メール / チャット / 電話）", C_START, fs=8)

    rbox(5, 8.6, 4.5, 0.65,
         "自動スコアリング＆重要度判定\n（AIクラシファイア）", C_PROC, fs=8)
    arrow(ax, 5, 9.17, 5, 8.92)

    diamond(5, 7.5, 4.0, 0.9, "重要度\n判定結果", C_DEC, fs=9)
    arrow(ax, 5, 8.27, 5, 7.95)

    # ── 3分岐 ────────────────────────────────────────────────────────────────
    rbox(1.5, 6.1, 2.6, 0.8,
         "緊急（S）\nL3エスカレーション\n30分以内に対応開始", C_CRIT, fs=7.5)
    arrow(ax, 3.0, 7.5, 2.8, 6.5, label='緊急(S)', lx=-0.6, ly=0.15)

    rbox(5, 6.1, 2.6, 0.8,
         "重要（A）\nL2対応\n2時間以内に対応開始", C_HIGH, fs=7.5)
    arrow(ax, 5, 7.05, 5, 6.5, label='重要(A)', lx=0.18)

    rbox(8.5, 6.1, 2.6, 0.8,
         "通常（B/C）\nL1対応\n翌営業日に対応開始", C_NORM, fs=7.5)
    arrow(ax, 7.0, 7.5, 7.2, 6.5, label='通常(B/C)', lx=0.1, ly=0.15)

    # ── 合流：調査・対応 ──────────────────────────────────────────────────────
    rbox(5, 4.8, 4.5, 0.65, "調査・対応・解決策の提案", C_PROC, fs=8.5)
    arrow(ax, 1.5, 5.7, 3.25, 5.12)
    arrow(ax, 5, 5.7, 5, 5.12)
    arrow(ax, 8.5, 5.7, 6.75, 5.12)

    rbox(5, 3.95, 4.5, 0.65, "顧客による解決確認", C_PROC, fs=8.5)
    arrow(ax, 5, 4.47, 5, 4.27)

    diamond(5, 3.05, 3.0, 0.75, "解決\n確認?", C_DEC, fs=9)
    arrow(ax, 5, 3.62, 5, 3.42)

    # Yes: クローズ
    rbox(5, 2.1, 3.0, 0.65, "チケットクローズ", C_END, fs=8.5)
    arrow(ax, 5, 2.67, 5, 2.42, label='Yes', lx=0.2)

    rbox(5, 1.2, 3.5, 0.65,
         "3日後フォローアップ調査\nCSATスコア取得", C_END, fs=8)
    arrow(ax, 5, 1.77, 5, 1.52)

    # No: ループバック
    ax.annotate('', xy=(3.5, 4.8), xytext=(3.5, 3.05),
                arrowprops=dict(arrowstyle='->', color=C_CRIT, lw=1.4,
                                connectionstyle='arc3,rad=-0.35'))
    ax.text(2.4, 3.9, 'No\n（再調査）', fontsize=8, color=C_CRIT,
            fontweight='bold', ha='center')

    ax.set_title('カスタマーサクセス 問い合わせ対応フロー（2024年版）',
                 fontsize=13, fontweight='bold', pad=8)

    buf_flow = io.BytesIO()
    fig.savefig(buf_flow, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    flow_bytes = buf_flow.getvalue()
    plt.close(fig)

    # ── DOCX 組み立て ────────────────────────────────────────────────────────
    def set_cell_bg(cell, hex_color):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_color)
        tcPr.append(shd)

    doc = DocxDocument()
    for section in doc.sections:
        section.page_height = Cm(29.7)
        section.page_width  = Cm(21.0)

    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_p.add_run("カスタマーサクセス 問い合わせ対応フロー")
    run.bold = True
    run.font.size = Pt(18)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.add_run("CS部 標準業務手順書　2024年4月版").font.size = Pt(12)

    doc.add_paragraph()

    doc.add_heading("1. 概要", level=2)
    doc.add_paragraph(
        "本フローは、顧客から寄せられた問い合わせを適切に分類・対応するための標準プロセスを定めたものです。"
        "AIによる自動スコアリングで重要度を4段階（S/A/B/C）に判定し、各区分に応じた対応体制を取ります。"
        "具体的な対応時間・エスカレーション条件・判定ロジックは下記フロー図を参照してください。"
    )

    doc.add_heading("2. 対応フロー図", level=2)
    doc.add_paragraph(
        "※ 重要度ごとの対応時間・担当チーム・再調査ループの条件は下図を参照してください。"
    )

    # フロー図PNG埋め込み
    pic_p = doc.add_paragraph()
    pic_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_pic = pic_p.add_run()
    run_pic.add_picture(io.BytesIO(flow_bytes), width=Inches(5.8))

    doc.add_paragraph()
    doc.add_heading("3. SLA・KPI目標値", level=2)

    sla_rows = [
        ("初回応答時間 （緊急S）",  "30分以内",   "99%以上"),
        ("初回応答時間 （重要A）",  "2時間以内",  "95%以上"),
        ("初回応答時間 （通常B/C）", "翌営業日",  "90%以上"),
        ("初回解決率",              "65%以上",    "─"),
        ("顧客満足度スコア（CSAT）", "4.0/5.0以上", "─"),
    ]
    tbl = doc.add_table(rows=len(sla_rows) + 1, cols=3)
    tbl.style = "Table Grid"
    for c_idx, h in enumerate(["SLA項目", "目標値", "達成率目標"]):
        cell = tbl.rows[0].cells[c_idx]
        cell.text = h
        set_cell_bg(cell, "1F3864")
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.paragraphs[0].runs[0].bold = True
    for r_idx, (item, target, rate) in enumerate(sla_rows, 1):
        tbl.rows[r_idx].cells[0].text = item
        tbl.rows[r_idx].cells[1].text = target
        tbl.rows[r_idx].cells[2].text = rate

    doc.save(DOCS / "業務フロー図_カスタマーサクセス対応.docx")
    print("  ✓ 業務フロー図_カスタマーサクセス対応.docx")


# ─── メイン ──────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=== 画像埋め込みドキュメント生成 ===")
    make_dashboard_pdf()
    make_org_chart_pptx()
    make_flowchart_docx()
    print("=== 完了 ===")
