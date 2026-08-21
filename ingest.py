"""
ドキュメントを読み込み、ChromaDBにベクトル化して保存するスクリプト。
対応形式: PDF, DOCX, PPTX, XLSX

変更履歴:
  v3: ローダーをlist[Document]返却に変更。PDF=ページ単位、PPTX=スライド単位、
      XLSX=シート単位、DOCX=見出し単位でDocument化。テーブルは独立Document
      (Markdown形式)としてsplitter対象外に。chunk_overlap 50→100。
      --collection オプションで複数バージョンのDBを管理可能。
"""

import argparse
import os
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document

import pymupdf
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl


DOCS_DIR = Path(__file__).parent / "docs"
CHROMA_DIR = Path(__file__).parent / "chroma_db"
EMBED_MODEL = "intfloat/multilingual-e5-small"
DEFAULT_COLLECTION = "dify_rag"


# ── ユーティリティ ────────────────────────────────────────────────────────────

def rows_to_markdown(rows: list[list[str]]) -> str:
    """行リストをMarkdownテーブル文字列に変換する。"""
    if not rows:
        return ""
    header = rows[0]
    lines = ["| " + " | ".join(header) + " |",
             "|" + "|".join(["---"] * len(header)) + "|"]
    for row in rows[1:]:
        # 列数を揃える
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[:len(header)]) + " |")
    return "\n".join(lines)


# ── ローダー（各形式 → list[Document]） ──────────────────────────────────────

def load_pdf(path: Path) -> list[Document]:
    """PDFをページ単位でDocument化する。
    PyMuPDFでテキスト抽出 + テーブル検出を行い、
    テーブルはMarkdown形式の独立Documentとして登録する。
    """
    doc = pymupdf.open(str(path))
    docs = []

    for page_num in range(len(doc)):
        page = doc[page_num]

        # ── テーブル検出（Markdown Document として独立登録） ──
        try:
            for i, table in enumerate(page.find_tables().tables):
                rows = table.extract()  # list[list[str|None]]
                if rows:
                    rows_clean = [
                        [str(c) if c is not None else "" for c in row]
                        for row in rows
                    ]
                    md = rows_to_markdown(rows_clean)
                    if md:
                        docs.append(Document(
                            page_content=f"[{path.name} p.{page_num + 1}]\n{md}",
                            metadata={
                                "source": path.name,
                                "file_type": "pdf",
                                "page": page_num + 1,
                                "type": "table",
                            },
                        ))
        except Exception:
            pass

        # ── テキスト抽出（タグ付きPDFは読み順が正確） ──
        text = page.get_text("text").strip()
        if text:
            docs.append(Document(
                page_content=text,
                metadata={"source": path.name, "file_type": "pdf", "page": page_num + 1},
            ))

    return docs


def load_docx(path: Path) -> list[Document]:
    """DOCXを見出し単位のセクションDocument + テーブルDocumentに分割する。"""
    doc = DocxDocument(str(path))
    docs = []

    # ── 見出し単位でセクション分割 ──
    current_heading = "（前文）"
    current_parts = []

    def flush_section():
        text = "\n".join(current_parts).strip()
        if text:
            docs.append(Document(
                page_content=text,
                metadata={"source": path.name, "file_type": "docx",
                          "section": current_heading},
            ))

    for p in doc.paragraphs:
        if p.style.name.startswith("Heading"):
            flush_section()
            current_heading = p.text.strip() or current_heading
            current_parts = [p.text] if p.text.strip() else []
        else:
            if p.text.strip():
                current_parts.append(p.text)
    flush_section()

    # ── テーブルを独立Document（Markdown）として追加 ──
    for i, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "｜") for cell in row.cells]
            rows.append(cells)
        md = rows_to_markdown(rows)
        if md:
            docs.append(Document(
                page_content=md,
                metadata={"source": path.name, "file_type": "docx",
                          "type": "table", "table_index": i + 1},
            ))

    # ── テキストボックス ──
    for elem in doc.element.body.iter():
        if elem.tag.endswith("}txbxContent"):
            txbx_text = "".join(
                t.text for t in elem.iter() if t.tag.endswith("}t") and t.text
            )
            if txbx_text.strip():
                docs.append(Document(
                    page_content=txbx_text,
                    metadata={"source": path.name, "file_type": "docx",
                              "type": "textbox"},
                ))

    return docs


def load_pptx(path: Path) -> list[Document]:
    """PPTXをスライド単位でDocument化する。テーブルは独立Document（Markdown）。"""
    prs = Presentation(str(path))
    docs = []

    for i, slide in enumerate(prs.slides):
        texts = []

        for shape in slide.shapes:
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("|", "｜") for cell in row.cells]
                    rows.append(cells)
                md = rows_to_markdown(rows)
                if md:
                    docs.append(Document(
                        page_content=md,
                        metadata={"source": path.name, "file_type": "pptx",
                                  "slide": i + 1, "type": "table"},
                    ))
            elif shape.has_chart:
                try:
                    chart = shape.chart
                    for plot in chart.plots:
                        cats = []
                        if hasattr(plot, "categories") and plot.categories:
                            cats = [str(c) if c is not None else "" for c in plot.categories]
                        for series in plot.series:
                            name = series.name or ""
                            vals = [str(v) if v is not None else "" for v in series.values]
                            if cats:
                                for cat, val in zip(cats, vals):
                                    if cat or val:
                                        texts.append(f"{name}\t{cat}\t{val}")
                            elif vals:
                                texts.append(f"{name}\t" + "\t".join(vals))
                except Exception:
                    pass
            elif hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)

        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={"source": path.name, "file_type": "pptx", "slide": i + 1},
            ))

    return docs


def load_xlsx(path: Path) -> list[Document]:
    """XLSXをシート単位でDocument化する（Markdownテーブル形式）。"""
    wb = openpyxl.load_workbook(str(path), data_only=True)
    docs = []

    for sheet in wb.worksheets:
        all_rows = list(sheet.iter_rows(values_only=True))
        non_empty = [
            [str(c) if c is not None else "" for c in row]
            for row in all_rows
            if any(c is not None for c in row)
        ]
        if not non_empty:
            continue

        header = f"[シート: {sheet.title}]\n"
        md = rows_to_markdown(non_empty)
        docs.append(Document(
            page_content=header + md,
            metadata={"source": path.name, "file_type": "xlsx",
                      "sheet": sheet.title, "type": "table"},
        ))

    return docs


LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_xlsx,
}


# ── メイン処理 ────────────────────────────────────────────────────────────────

def load_documents() -> list[Document]:
    documents = []
    for path in sorted(DOCS_DIR.iterdir()):
        ext = path.suffix.lower()
        loader = LOADERS.get(ext)
        if loader is None:
            print(f"  スキップ（未対応形式）: {path.name}")
            continue
        print(f"  読み込み中: {path.name}")
        try:
            docs = loader(path)
            documents.extend(docs)
        except Exception as e:
            print(f"  エラー ({path.name}): {e}")
    return documents


def main():
    parser = argparse.ArgumentParser(description="ドキュメントをChromaDBに登録する")
    parser.add_argument(
        "--collection",
        default=DEFAULT_COLLECTION,
        help=f"ChromaDBのコレクション名 (default: {DEFAULT_COLLECTION})",
    )
    args = parser.parse_args()

    print(f"=== ドキュメント読み込み開始 ===")
    raw_docs = load_documents()
    print(f"読み込み完了: {len(raw_docs)} Documents\n")

    # テーブルDocumentはsplitter対象外、テキストDocumentのみ分割
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100,
        separators=["\n\n", "\n", "。", "、", " ", ""],
    )

    chunks = []
    table_count = 0
    for doc in raw_docs:
        if doc.metadata.get("type") == "table":
            chunks.append(doc)  # テーブルは分割しない
            table_count += 1
        else:
            chunks.extend(splitter.split_documents([doc]))

    print(f"=== テキスト分割 ===")
    print(f"  テーブルDocument（分割なし）: {table_count}件")
    print(f"  テキストチャンク（分割後）: {len(chunks) - table_count}件")
    print(f"  合計: {len(chunks)}件\n")

    print(f"=== 埋め込みモデル読み込み: {EMBED_MODEL} ===")
    embeddings = HuggingFaceEmbeddings(
        model_name=EMBED_MODEL,
        model_kwargs={"device": "cpu"},
        encode_kwargs={"normalize_embeddings": True},
    )

    print(f"=== ChromaDB に保存: {CHROMA_DIR} / collection={args.collection} ===")
    vectorstore = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=str(CHROMA_DIR),
        collection_name=args.collection,
    )
    print(f"保存完了: {vectorstore._collection.count()} チャンクを登録しました。")


if __name__ == "__main__":
    main()
