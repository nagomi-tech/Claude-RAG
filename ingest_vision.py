"""
Vision-based image indexing for all documents under docs/.

Scans docs/ for PDF/PPTX/DOCX/XLSX files and extracts all embedded images.
Each image is described by Claude Vision and added to ChromaDB as a text chunk.

Usage:
  python ingest_vision.py                              # process all files
  python ingest_vision.py --force                      # re-index even if already indexed
  python ingest_vision.py --collection dify_rag_v3     # 指定コレクションに登録

# ── 既知の制限事項（v2.1時点） ──────────────────────────────────────────────
#
# 【PPTX】グループシェイプ内の画像が未対応
#   PowerPoint でグループ化されたオブジェクト（GROUP_SHAPE）の中に画像がある場合、
#   再帰的な探索を行っていないため画像がスキップされる。
#
# 【DOCX】浮動配置の画像が未対応
#   python-docx の doc.inline_shapes はインライン配置（文字列中に埋め込まれた画像）のみ対象。
#   「文字列の折り返し」を設定した浮動画像（Word の「前面」「四角形」等）は取得されない。
#
# 【全形式】画像のメディアタイプ判定がPNG・JPEGのみ
#   先頭バイトがPNGシグネチャ以外の画像はすべて image/jpeg として Claude API に送信する。
#   GIF・WebP・BMP 等が含まれると誤ったメディアタイプとなり、APIエラーになる可能性がある。
#   Claude API がサポートするメディアタイプ: image/jpeg, image/png, image/gif, image/webp
#
# ─────────────────────────────────────────────────────────────────────────────
"""

import argparse
import base64
import io
import os
import time
from pathlib import Path

import anthropic
from dotenv import load_dotenv
from langchain_chroma import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document

load_dotenv()

DOCS = Path(__file__).parent / "docs"
CHROMA_DIR = Path(__file__).parent / "chroma_db"
EMBED_MODEL = "intfloat/multilingual-e5-small"
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")


# ─── 画像抽出（ファイル形式別） ───────────────────────────────────────────────

def pdf_images(pdf_path: Path) -> list[tuple[str, bytes]]:
    """PDFからVision対象の画像を抽出する。

    ページ種別を自動判定して処理を振り分ける:
      - ラスタ埋め込み画像あり → XObjectを直接抽出してVisionへ
      - テキスト少（<100文字）かつラスタなし → ページ全体をレンダリングしてVisionへ
        （ベクターグラフ・スキャンPDF対応）
      - テキスト十分 → Vision不要（ingest.pyのテキスト抽出でカバー）
    """
    import pymupdf

    TEXT_THRESHOLD = 100  # これ未満の文字数を「テキスト不十分」と判断
    doc = pymupdf.open(str(pdf_path))
    results = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text").strip()

        # ── ラスタ埋め込み画像を抽出 ──
        image_list = page.get_images(full=True)
        raster_found = False
        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            try:
                base_img = doc.extract_image(xref)
                img_bytes = base_img["image"]
                results.append((f"p.{page_num + 1}-図{img_idx + 1}", img_bytes))
                raster_found = True
            except Exception:
                pass

        # ── テキスト不十分 + ラスタなし → ページレンダリング ──
        # ベクターグラフ（Word/PowerPointのネイティブグラフ→PDF出力）や
        # スキャンPDFに対応するフォールバック
        if len(text) < TEXT_THRESHOLD and not raster_found:
            mat = pymupdf.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            results.append((f"p.{page_num + 1}-レンダリング", img_bytes))

    return results


def pptx_images(pptx_path: Path) -> list[tuple[str, bytes]]:
    """Extract Picture shapes from all PPTX slides."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(pptx_path))
    results = []
    for i, slide in enumerate(prs.slides, 1):
        for j, shape in enumerate(slide.shapes, 1):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                results.append((f"スライド{i}-図{j}", shape.image.blob))
    return results


def docx_images(docx_path: Path) -> list[tuple[str, bytes]]:
    """Extract inline images from DOCX."""
    from docx import Document as DocxDocument

    doc = DocxDocument(str(docx_path))
    results = []
    for i, shape in enumerate(doc.inline_shapes, 1):
        try:
            rId = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            blob = doc.part.related_parts[rId].blob
            results.append((f"図{i}", blob))
        except Exception:
            pass
    return results


def xlsx_images(xlsx_path: Path) -> list[tuple[str, bytes]]:
    """Extract embedded image objects from all worksheets in an XLSX."""
    from openpyxl import load_workbook

    wb = load_workbook(str(xlsx_path), data_only=True)
    results = []
    for ws in wb.worksheets:
        for i, img in enumerate(ws._images, 1):
            try:
                data = img._data()
                results.append((f"{ws.title}-画像{i}", data))
            except Exception:
                pass
    return results


def extract_images(path: Path) -> list[tuple[str, bytes]]:
    """Dispatch to the appropriate extractor by file extension."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return pdf_images(path)
        elif ext in (".pptx", ".ppt"):
            return pptx_images(path)
        elif ext in (".docx", ".doc"):
            return docx_images(path)
        elif ext in (".xlsx", ".xls"):
            return xlsx_images(path)
    except Exception as e:
        print(f"    [抽出エラー] {path.name}: {e}")
    return []


# ─── Vision 説明文生成 ────────────────────────────────────────────────────────

def describe_image(client: anthropic.Anthropic, image_bytes: bytes,
                   doc_name: str, image_label: str) -> str:
    """Send image to Claude Vision and return Japanese text description."""
    b64 = base64.standard_b64encode(image_bytes).decode()
    media_type = "image/png" if image_bytes[:4] == b"\x89PNG" else "image/jpeg"

    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=2048,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image",
                 "source": {"type": "base64", "media_type": media_type, "data": b64}},
                {"type": "text",
                 "text": (
                     f"これは「{doc_name}」の{image_label}です。\n"
                     "この画像に含まれるすべての情報（数値、名前、項目名、構造、関係性など）を"
                     "漏れなく日本語のテキストとして書き起こしてください。\n"
                     "グラフの場合は各項目の値を、組織図の場合は階層構造と人名を、"
                     "フロー図の場合は各ステップと条件分岐を具体的に記述してください。"
                 )},
            ],
        }],
    )
    return msg.content[0].text


# ─── ChromaDB ────────────────────────────────────────────────────────────────

def get_vectorstore(collection: str) -> Chroma:
    embeddings = HuggingFaceEmbeddings(
        model_name=EMBED_MODEL,
        model_kwargs={"device": "cpu"},
        encode_kwargs={"normalize_embeddings": True},
    )
    return Chroma(
        persist_directory=str(CHROMA_DIR),
        embedding_function=embeddings,
        collection_name=collection,
    )


def already_indexed(vectorstore: Chroma, source: str) -> bool:
    """Return True if vision_extracted chunks for this source already exist."""
    result = vectorstore.get(where={"$and": [
        {"source": {"$eq": source}},
        {"file_type": {"$eq": "vision_extracted"}},
    ]})
    return len(result["ids"]) > 0


# ─── メイン ───────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--force", action="store_true",
                        help="Re-index even if already indexed")
    parser.add_argument("--collection", default="dify_rag",
                        help="ChromaDBのコレクション名 (default: dify_rag)")
    args = parser.parse_args()

    if not ANTHROPIC_API_KEY:
        raise SystemExit("Error: ANTHROPIC_API_KEY not set in .env")

    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    vectorstore = get_vectorstore(args.collection)

    # docs/ 以下の対象ファイルをすべてスキャン
    target_extensions = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".xlsx", ".xls"}
    doc_files = sorted(
        f for f in DOCS.iterdir()
        if f.is_file() and f.suffix.lower() in target_extensions
    )

    print(f"対象ファイル: {len(doc_files)} 件\n")

    tasks: list[tuple[str, str, bytes]] = []  # (source, label, image_bytes)

    for doc_path in doc_files:
        source = doc_path.name
        if not args.force and already_indexed(vectorstore, source):
            print(f"  スキップ（既存）: {source}")
            continue

        images = extract_images(doc_path)
        if not images:
            print(f"  画像なし: {source}")
            continue

        print(f"  {source}: {len(images)} 画像")
        for label, img_bytes in images:
            tasks.append((source, label, img_bytes))

    if not tasks:
        print("\n新たにインデックスする画像はありません。(--force で再実行)")
        return

    print(f"\n合計 {len(tasks)} 件の画像をVisionで解析します...\n")

    docs_to_add: list[Document] = []
    for source, label, img_bytes in tasks:
        print(f"  [{source}] {label}", end=" ... ", flush=True)
        try:
            text = describe_image(client, img_bytes, source, label)
            docs_to_add.append(Document(
                page_content=text,
                metadata={
                    "source": source,
                    "file_type": "vision_extracted",
                    "image_desc": label,
                },
            ))
            print(f"{len(text)}文字")
        except Exception as e:
            print(f"[Visionエラー] {e}")
        time.sleep(1.0)

    if docs_to_add:
        vectorstore.add_documents(docs_to_add)
        print(f"\nChromaDBに {len(docs_to_add)} チャンク追加完了。")
        print("api.py を再起動してください。")


if __name__ == "__main__":
    main()
