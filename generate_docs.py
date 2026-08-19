"""
評価用ダミーデータ書類を10種類生成するスクリプト。
DOCX x3 / PPTX x3 / XLSX x2 / PDF x2
"""

import io
import random
from pathlib import Path

DOCS = Path(__file__).parent / "docs"
DOCS.mkdir(exist_ok=True)

# ─────────────────────────────────────────────
#  DOCX
# ─────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement


def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)


def add_heading(doc, text, level=1, color="1F3864"):
    h = doc.add_heading(text, level=level)
    run = h.runs[0] if h.runs else h.add_run(text)
    run.font.color.rgb = RGBColor.from_string(color)
    return h


# ── DOCX 1: 取締役会議事録 ──────────────────
def make_board_minutes():
    doc = Document()
    doc.core_properties.title = "取締役会議事録"

    # 用紙設定
    for section in doc.sections:
        section.page_height = Cm(29.7)
        section.page_width  = Cm(21.0)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("第18期 第3回 取締役会議事録")
    run.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    doc.add_paragraph()

    # 基本情報テーブル
    info_table = doc.add_table(rows=5, cols=4)
    info_table.style = "Table Grid"
    info_data = [
        ("開催日時", "2024年3月15日（金）14:00〜16:30", "開催場所", "本社 第1会議室（オンライン併用）"),
        ("議長",   "代表取締役社長　田中 一郎",         "記録者",   "総務部　山田 花子"),
        ("取締役数","8名（定員8名）",                    "出席取締役","8名（全員出席）",),
        ("監査役数","3名（定員3名）",                    "出席監査役","3名（全員出席）",),
        ("決議要件","取締役の過半数の出席かつ出席取締役の過半数の賛成", "成立要件","充足（可決要件：出席者5名以上の賛成）"),
    ]
    for r, row_data in enumerate(info_data):
        cells = info_table.rows[r].cells
        for i, text in enumerate(row_data):
            cells[i].text = text
            if i % 2 == 0:
                set_cell_bg(cells[i], "D6E4F0")
                cells[i].paragraphs[0].runs[0].bold = True

    doc.add_paragraph()
    add_heading(doc, "1. 出席者一覧", level=2)

    # 出席者テーブル
    attendees = [
        ["役職",       "氏名",       "区分",   "備考"],
        ["代表取締役社長", "田中 一郎", "取締役", "議長"],
        ["専務取締役（CFO）", "鈴木 二郎", "取締役", ""],
        ["常務取締役（CTO）", "佐藤 三郎", "取締役", ""],
        ["取締役（営業本部長）", "高橋 四郎", "取締役", ""],
        ["取締役（生産本部長）", "渡辺 五郎", "取締役", ""],
        ["社外取締役", "伊藤 花子", "取締役", "監査委員"],
        ["社外取締役", "加藤 太郎", "取締役", "報酬委員"],
        ["社外取締役", "山本 次郎", "取締役", "指名委員"],
        ["常勤監査役", "松本 一子", "監査役", ""],
        ["社外監査役", "井上 聡",   "監査役", ""],
        ["社外監査役", "木村 誠",   "監査役", ""],
    ]
    atbl = doc.add_table(rows=len(attendees), cols=4)
    atbl.style = "Table Grid"
    for r, row in enumerate(attendees):
        cells = atbl.rows[r].cells
        for c, val in enumerate(row):
            cells[c].text = val
            if r == 0:
                set_cell_bg(cells[c], "1F3864")
                p = cells[c].paragraphs[0]
                p.runs[0].bold = True
                p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    doc.add_paragraph()
    add_heading(doc, "2. 審議事項・決議内容", level=2)

    agenda_items = [
        ("第1号議案", "2023年度（第18期）決算の承認",
         "CFOより第18期（2023年4月〜2024年3月）の連結決算概要について報告があった。\n"
         "売上高：128億円（前期比+8.4%）、営業利益：15.6億円（前期比+12.1%）、\n"
         "純利益：10.2億円（前期比+9.7%）。各種指標について質疑応答の後、全会一致で承認。",
         "承認（賛成8名 / 反対0名 / 棄権0名）"),
        ("第2号議案", "2024年度事業計画・予算の承認",
         "営業本部長より2024年度の事業計画概要について説明があった。売上高目標140億円（前期比+9.4%）、"
         "新規事業（AI活用サービス）への投資15億円を含む設備投資計画について審議。"
         "投資規模の適切性について複数の質問が出され、CTOより詳細な投資効果説明があった後、賛多数で承認。",
         "承認（賛成7名 / 反対0名 / 棄権1名）"),
        ("第3号議案", "役員報酬総額の改定",
         "報酬委員長より現行の役員報酬体系の見直し案について説明があった。"
         "業績連動部分の比率を現行30%から40%に引き上げる案について審議。"
         "グローバルベンチマークとの比較データを参照しながら議論を行い、全会一致で承認。",
         "承認（賛成8名 / 反対0名 / 棄権0名）"),
        ("第4号議案", "新規子会社設立の承認",
         "CTOより AI 技術開発を専業とする完全子会社「株式会社テックラボ」（仮称）の設立計画について説明。"
         "資本金3億円、従業員数50名（初年度）の規模で2024年7月設立予定。"
         "事業計画・財務計画・リスク分析について質疑応答の後、全会一致で承認。",
         "承認（賛成8名 / 反対0名 / 棄権0名）"),
    ]

    for num, title, body, result in agenda_items:
        add_heading(doc, f"{num}：{title}", level=3)
        p = doc.add_paragraph(body)
        p.paragraph_format.left_indent = Cm(0.5)

        result_tbl = doc.add_table(rows=1, cols=2)
        result_tbl.style = "Table Grid"
        cells = result_tbl.rows[0].cells
        cells[0].text = "決議結果"
        set_cell_bg(cells[0], "1F3864")
        cells[0].paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cells[0].paragraphs[0].runs[0].bold = True
        cells[1].text = result
        doc.add_paragraph()

    add_heading(doc, "3. 報告事項", level=2)
    reports = [
        ("月次業績報告（2024年2月）", "売上高11.2億円（予算比+3.2%）。主要顧客向け案件の進捗良好。"),
        ("コンプライアンス報告", "2024年1〜2月の法令違反ゼロ。内部通報1件対応済み。"),
        ("IT セキュリティ報告", "インシデント0件。ランサムウェア対策強化のため追加投資を検討中。"),
    ]
    for rtitle, rbody in reports:
        add_heading(doc, f"■ {rtitle}", level=3)
        doc.add_paragraph(rbody)

    doc.add_paragraph()
    close = doc.add_paragraph("以上をもって本日の審議事項をすべて終了し、議長が閉会を宣言した。\n\n"
                              "以上、本議事録の正確性を証するため議長及び出席取締役が署名する。")
    sign_tbl = doc.add_table(rows=3, cols=2)
    sign_tbl.style = "Table Grid"
    signs = [("議長（代表取締役社長）", "田中 一郎　　　　㊞"),
             ("常務取締役（CTO）",     "佐藤 三郎　　　　㊞"),
             ("社外取締役",            "伊藤 花子　　　　㊞")]
    for r, (role, name) in enumerate(signs):
        sign_tbl.rows[r].cells[0].text = role
        sign_tbl.rows[r].cells[1].text = name

    doc.save(DOCS / "取締役会議事録_2024年3月期.docx")
    print("  ✓ 取締役会議事録_2024年3月期.docx")


# ── DOCX 2: 業務マニュアル ────────────────────
def make_manual():
    doc = Document()
    doc.core_properties.title = "受注処理業務マニュアル"

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("受注処理業務マニュアル　第5版")
    run.bold = True
    run.font.size = Pt(20)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run("営業管理部　2024年4月1日改訂").font.color.rgb = RGBColor(0x70, 0x70, 0x70)

    doc.add_paragraph()

    # 改訂履歴
    add_heading(doc, "改訂履歴", level=2)
    hist_tbl = doc.add_table(rows=6, cols=4)
    hist_tbl.style = "Table Grid"
    hist_header = ["改訂版", "改訂日", "改訂者", "主な改訂内容"]
    for c, h in enumerate(hist_header):
        cell = hist_tbl.rows[0].cells[c]
        cell.text = h
        set_cell_bg(cell, "2E74B5")
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        cell.paragraphs[0].runs[0].bold = True
    hist_data = [
        ("第1版","2020年4月1日","営業管理部","初版作成"),
        ("第2版","2021年10月1日","営業管理部","ECサイト連携フロー追加"),
        ("第3版","2022年7月1日","営業管理部","与信審査プロセス強化"),
        ("第4版","2023年4月1日","営業管理部","電子契約対応・押印廃止"),
        ("第5版","2024年4月1日","営業管理部","AI自動審査システム導入対応"),
    ]
    for r, row in enumerate(hist_data, 1):
        for c, val in enumerate(row):
            hist_tbl.rows[r].cells[c].text = val

    doc.add_paragraph()
    add_heading(doc, "1. 目的と適用範囲", level=2)
    doc.add_paragraph(
        "本マニュアルは、受注処理業務の標準化・品質向上を目的として作成した。"
        "国内法人顧客からの新規注文受付から出荷指示までのすべての業務フローに適用する。"
        "個人顧客向けEC受注については別マニュアル（EC受注処理規程）を参照すること。"
    )

    add_heading(doc, "2. 業務フロー概要", level=2)
    flow_steps = [
        ("STEP 1", "注文受付", "FAX・メール・EDI・Web受注フォームにて受付。\nシステムへ24時間以内に入力すること。"),
        ("STEP 2", "与信確認", "AI与信審査システム（AutoCredit v3）で自動判定。\nスコア70以上：自動承認 / 50〜69：担当者確認 / 49以下：要申請。"),
        ("STEP 3", "在庫確認", "基幹システム（ERP）で在庫照会。在庫不足の場合は生産管理部へ増産依頼フォームを送付。"),
        ("STEP 4", "受注確定", "与信OK・在庫確認完了後、顧客へ受注確認書をメール送付（自動送信）。"),
        ("STEP 5", "出荷指示", "出荷日3営業日前までに倉庫管理システムへ出荷指示を登録。伝票・梱包仕様を添付。"),
        ("STEP 6", "完了確認", "出荷完了後、ERPの出荷ステータスをCOMPLETEに更新。請求書は翌月末発行。"),
    ]
    flow_tbl = doc.add_table(rows=len(flow_steps)+1, cols=3)
    flow_tbl.style = "Table Grid"
    for c, h in enumerate(["ステップ", "処理名", "内容・注意事項"]):
        cell = flow_tbl.rows[0].cells[c]
        cell.text = h
        set_cell_bg(cell, "1F3864")
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        cell.paragraphs[0].runs[0].bold = True
    for r, (step, name, desc) in enumerate(flow_steps, 1):
        flow_tbl.rows[r].cells[0].text = step
        flow_tbl.rows[r].cells[1].text = name
        flow_tbl.rows[r].cells[2].text = desc
        if r % 2 == 0:
            for c in range(3):
                set_cell_bg(flow_tbl.rows[r].cells[c], "EBF3FB")

    doc.add_paragraph()
    add_heading(doc, "3. 与信審査基準", level=2)
    criteria_tbl = doc.add_table(rows=6, cols=4)
    criteria_tbl.style = "Table Grid"
    criteria_header = ["判定区分", "スコア範囲", "処理方法", "対応期限"]
    for c, h in enumerate(criteria_header):
        cell = criteria_tbl.rows[0].cells[c]
        cell.text = h
        set_cell_bg(cell, "375623")
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        cell.paragraphs[0].runs[0].bold = True
    criteria_data = [
        ("自動承認",   "70〜100点", "システム自動承認・即時受注確定",          "0時間（即時）"),
        ("担当者確認", "60〜69点",  "営業担当者が顧客情報を確認後、判断",      "当日中"),
        ("上長承認",   "50〜59点",  "部長クラスの承認が必要",                  "翌営業日まで"),
        ("与信申請",   "30〜49点",  "財務部へ与信枠拡大申請書を提出",          "3営業日"),
        ("受注停止",   "0〜29点",   "受注不可。顧客へ理由を説明し代替案を提示", "即時連絡"),
    ]
    for r, row in enumerate(criteria_data, 1):
        for c, val in enumerate(row):
            criteria_tbl.rows[r].cells[c].text = val

    add_heading(doc, "4. よくあるトラブルと対応", level=2)
    troubles = [
        ("在庫切れ時の対応",
         "① 生産管理部へ緊急増産依頼フォーム（Form-P03）を送付\n"
         "② 顧客へ納期延長の連絡（納期+5営業日が目安）\n"
         "③ 代替品がある場合は顧客に提案。承諾得た場合は別注文として処理"),
        ("EDI受信エラー",
         "① IT部門へエラーコードを報告\n"
         "② 顧客へFAXにて注文書の再送を依頼\n"
         "③ 手動入力後、備考欄に「EDIエラーによる手動入力」と記載"),
        ("与信審査の長期未完了",
         "① 3営業日経過後も審査未完了の場合、財務部長へエスカレーション\n"
         "② 顧客への状況報告を2日ごとに実施\n"
         "③ 最長10営業日で可否を決定すること"),
    ]
    for title, desc in troubles:
        add_heading(doc, f"■ {title}", level=3)
        doc.add_paragraph(desc)

    doc.save(DOCS / "業務マニュアル_受注処理フロー_第5版.docx")
    print("  ✓ 業務マニュアル_受注処理フロー_第5版.docx")


# ── DOCX 3: 採用計画書 ───────────────────────
def make_recruitment():
    doc = Document()
    doc.core_properties.title = "2024年度採用計画書"

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("2024年度　採用計画書（新卒・中途）")
    run.bold = True
    run.font.size = Pt(18)

    doc.add_paragraph()
    add_heading(doc, "1. 採用方針", level=2)
    doc.add_paragraph(
        "2024年度は事業拡大フェーズに対応するため、エンジニア・営業・コーポレートの三領域で"
        "積極的な採用を実施する。特にAI・データサイエンス領域の即戦力中途採用に注力し、"
        "新卒採用においても技術系人材の拡充を図る。"
    )

    add_heading(doc, "2. 採用計画数（職種別）", level=2)
    plan_tbl = doc.add_table(rows=11, cols=6)
    plan_tbl.style = "Table Grid"
    plan_header = ["部門", "職種", "新卒", "中途", "合計", "優先度"]
    for c, h in enumerate(plan_header):
        cell = plan_tbl.rows[0].cells[c]
        cell.text = h
        set_cell_bg(cell, "C00000")
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        cell.paragraphs[0].runs[0].bold = True
    plan_data = [
        ("技術本部", "ソフトウェアエンジニア（バックエンド）", "5", "10", "15", "★★★"),
        ("技術本部", "AIエンジニア・MLエンジニア",             "3", "8",  "11", "★★★"),
        ("技術本部", "フロントエンドエンジニア",               "3", "5",  "8",  "★★☆"),
        ("技術本部", "インフラ・SRE",                          "1", "4",  "5",  "★★☆"),
        ("技術本部", "QA・品質管理",                           "2", "3",  "5",  "★☆☆"),
        ("営業本部", "法人営業（フィールドセールス）",          "5", "7",  "12", "★★★"),
        ("営業本部", "インサイドセールス",                      "3", "3",  "6",  "★★☆"),
        ("マーケ",  "デジタルマーケター",                       "2", "3",  "5",  "★★☆"),
        ("コーポ",  "人事・採用担当",                           "1", "2",  "3",  "★☆☆"),
        ("コーポ",  "経理・財務",                               "1", "2",  "3",  "★☆☆"),
    ]
    for r, row in enumerate(plan_data, 1):
        for c, val in enumerate(row):
            plan_tbl.rows[r].cells[c].text = val
        if r % 2 == 0:
            for c in range(6):
                set_cell_bg(plan_tbl.rows[r].cells[c], "FFF0F0")

    doc.add_paragraph()
    add_heading(doc, "3. 採用チャネル別コスト計画", level=2)
    cost_tbl = doc.add_table(rows=8, cols=5)
    cost_tbl.style = "Table Grid"
    cost_header = ["採用チャネル", "活用対象", "予算（万円）", "目標採用数", "1人当たり単価（万円）"]
    for c, h in enumerate(cost_header):
        cell = cost_tbl.rows[0].cells[c]
        cell.text = h
        set_cell_bg(cell, "1F3864")
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        cell.paragraphs[0].runs[0].bold = True
    cost_data = [
        ("新卒：就活ナビ",     "新卒全職種",   "1,200", "26", "46"),
        ("新卒：学内説明会",   "新卒理系",     "300",   "8",  "38"),
        ("中途：人材紹介",     "中途エンジニア","6,000","30", "200"),
        ("中途：求人媒体",     "中途全職種",   "1,800", "15", "120"),
        ("リファラル採用",     "全職種",       "500",   "10", "50"),
        ("SNS採用（LinkedIn）","エンジニア・営業","200", "5", "40"),
        ("自社採用サイト",     "全職種",       "100",   "3",  "33"),
    ]
    for r, row in enumerate(cost_data, 1):
        for c, val in enumerate(row):
            cost_tbl.rows[r].cells[c].text = val

    add_heading(doc, "4. 採用スケジュール（新卒）", level=2)
    sched_tbl = doc.add_table(rows=8, cols=7)
    sched_tbl.style = "Table Grid"
    months = ["施策", "3月", "4月", "5月", "6月", "7月", "8月"]
    for c, h in enumerate(months):
        cell = sched_tbl.rows[0].cells[c]
        cell.text = h
        set_cell_bg(cell, "203864")
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        cell.paragraphs[0].runs[0].bold = True
    sched_data = [
        ("会社説明会",    "●", "●", "●", "",  "",  ""),
        ("ES受付",        "●", "●", "●", "",  "",  ""),
        ("1次選考",       "",  "●", "●", "●", "",  ""),
        ("2次選考",       "",  "",  "●", "●", "●", ""),
        ("最終面接",      "",  "",  "",  "●", "●", ""),
        ("内定通知",      "",  "",  "",  "●", "●", "●"),
        ("内定者フォロー","",  "",  "",  "●", "●", "●"),
    ]
    for r, row in enumerate(sched_data, 1):
        sched_tbl.rows[r].cells[0].text = row[0]
        for c in range(1, 7):
            sched_tbl.rows[r].cells[c].text = row[c]
            if row[c] == "●":
                set_cell_bg(sched_tbl.rows[r].cells[c], "BFEFBF")

    doc.save(DOCS / "2024年度採用計画書.docx")
    print("  ✓ 2024年度採用計画書.docx")


# ─────────────────────────────────────────────
#  PPTX
# ─────────────────────────────────────────────
from pptx import Presentation as PRS
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm as PCm
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor as PptxColor


def add_slide_title(prs, layout_idx, title_text, subtitle_text=""):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    tf = slide.shapes.title
    if tf:
        tf.text = title_text
    if subtitle_text and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = subtitle_text
        except:
            pass
    return slide


def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=(0,0,0), align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = PptxColor(*color)
    return tb


# ── PPTX 1: Q1事業実績報告 ──────────────────
def make_q1_report():
    prs = PRS()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # スライド1: タイトル
    s = add_slide_title(prs, 0, "2024年度 第1四半期 事業実績報告",
                        "2024年7月10日　取締役会資料")

    # スライド2: サマリ KPI
    s2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_text_box(s2, "エグゼクティブサマリー　─ 2024Q1 業績ハイライト",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=22, bold=True, color=(31,56,100))

    kpis = [
        ("売上高", "35.2億円", "+12.4%", "目標比+2.1%"),
        ("営業利益", "4.8億円", "+18.7%", "目標比+0.8%"),
        ("営業利益率", "13.6%", "+0.7pt", "目標差±0pt"),
        ("受注高", "41.5億円", "+9.2%", "目標比+3.5%"),
    ]
    box_w = Inches(2.9)
    for i, (label, val, yoy, vs_tgt) in enumerate(kpis):
        x = Inches(0.4 + i * 3.1)
        # 背景ボックス（白）
        box = s2.shapes.add_shape(1, x, Inches(1.0), box_w, Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = PptxColor(240, 246, 255)
        box.line.color.rgb = PptxColor(31, 56, 100)

        add_text_box(s2, label, x+Inches(0.1), Inches(1.1), Inches(2.7), Inches(0.4),
                     font_size=13, bold=True, color=(31,56,100))
        add_text_box(s2, val,   x+Inches(0.1), Inches(1.5), Inches(2.7), Inches(0.6),
                     font_size=24, bold=True, color=(0,70,127))
        add_text_box(s2, f"前年同期比 {yoy}", x+Inches(0.1), Inches(2.1), Inches(2.7), Inches(0.35),
                     font_size=12, color=(0,128,0))
        add_text_box(s2, vs_tgt, x+Inches(0.1), Inches(2.45), Inches(2.7), Inches(0.35),
                     font_size=11, color=(100,100,100))

    add_text_box(s2, "▼ 主要トピック\n"
                     "① AI 営業支援ツール導入により既存顧客への提案件数が前四半期比+34%\n"
                     "② 新製品「ProSeries X」の受注が計画の150%で好調なスタート\n"
                     "③ 原材料費上昇（前年比+8%）の影響を生産効率化でカバー。COGS率改善",
                 Inches(0.5), Inches(3.5), Inches(12), Inches(2.5),
                 font_size=13, color=(30,30,30))

    # スライド3: 月別売上グラフ
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s3, "月次売上高・営業利益の推移（2024年4〜6月）",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=20, bold=True, color=(31,56,100))

    chart_data = CategoryChartData()
    chart_data.categories = ["4月", "5月", "6月"]
    chart_data.add_series("売上高（億円）",  (10.8, 11.5, 12.9))
    chart_data.add_series("営業利益（億円）", (1.3,  1.6,  1.9))

    chart = s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.0), Inches(7.5), Inches(5.5),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "月別売上・利益（億円）"

    add_text_box(s3, "【分析コメント】\n"
                     "・4月：期初の立ち上がりは想定通り。新年度予算の顧客稟議完了待ちの案件多数。\n"
                     "・5月：大型案件（A社 3.5億）の受注完了でジャンプアップ。\n"
                     "・6月：四半期末効果もあり最高売上。利益率も改善し通期目標に対し順調。",
                 Inches(8.2), Inches(1.0), Inches(4.8), Inches(5.5),
                 font_size=12, color=(50,50,50))

    # スライド4: 部門別売上
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s4, "部門別売上構成と前年同期比較",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=20, bold=True, color=(31,56,100))

    pie_data = ChartData()
    pie_data.categories = ["製品事業", "サービス事業", "保守・サポート", "新規事業"]
    pie_data.add_series("Q1売上", (18.2, 9.5, 5.8, 1.7))

    pie = s4.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.3), Inches(1.0), Inches(6.0), Inches(5.5),
        pie_data,
    ).chart
    pie.has_title = True
    pie.chart_title.text_frame.text = "部門別売上構成（2024Q1）"

    dept_tbl_data = [
        ["部門",       "Q1実績", "前年Q1", "前年比",    "通期計画"],
        ["製品事業",   "18.2億", "15.8億", "+15.2%",   "72.0億"],
        ["サービス事業", "9.5億", "8.9億",  "+6.7%",    "38.0億"],
        ["保守・サポート", "5.8億", "5.6億", "+3.6%",   "23.0億"],
        ["新規事業",   "1.7億",  "0.6億",  "+183.3%",  "7.0億"],
        ["合計",       "35.2億", "30.9億", "+13.9%",   "140.0億"],
    ]
    dept_x, dept_y = Inches(6.5), Inches(1.5)
    from pptx.util import Pt as PPt
    tbl = s4.shapes.add_table(len(dept_tbl_data), 5,
                               dept_x, dept_y, Inches(6.5), Inches(4.5)).table
    for r, row in enumerate(dept_tbl_data):
        for c, val in enumerate(row):
            tbl.cell(r,c).text = val
            if r == 0:
                tbl.cell(r,c).fill.solid()
                tbl.cell(r,c).fill.fore_color.rgb = PptxColor(31,56,100)
                tbl.cell(r,c).text_frame.paragraphs[0].runs[0].font.color.rgb = PptxColor(255,255,255)
                tbl.cell(r,c).text_frame.paragraphs[0].runs[0].font.bold = True

    # スライド5: 課題と次四半期予測
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s5, "課題・リスクと第2四半期見通し",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=20, bold=True, color=(31,56,100))
    add_text_box(s5, "【主要課題・リスク】\n"
                     "① 半導体不足による製品製造リードタイム延長（影響額：最大 2億円）\n"
                     "② 競合他社の新製品投入による価格競争激化リスク（国内市場）\n"
                     "③ 為替影響：円安進行により輸入原材料費が増加傾向（Q2影響 +0.5億円見込み）\n\n"
                     "【Q2 売上高・営業利益予測】\n"
                     "売上高目標：37.5億円（+6.5% vs Q1）\n"
                     "営業利益目標：5.3億円（+10.4% vs Q1）",
                 Inches(0.5), Inches(1.0), Inches(12), Inches(5.5),
                 font_size=14, color=(40,40,40))

    prs.save(DOCS / "2024年Q1事業実績報告書.pptx")
    print("  ✓ 2024年Q1事業実績報告書.pptx")


# ── PPTX 2: 新製品企画提案書 ────────────────
def make_product_proposal():
    prs = PRS()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_slide_title(prs, 0, "AI営業支援プラットフォーム「SalesAI Pro」\n企画提案書",
                    "2024年5月　プロダクト戦略部")

    # スライド2: 背景と課題
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s2, "市場背景と顧客課題", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=22, bold=True, color=(31,56,100))
    problems = [
        ("営業生産性の低下", "営業担当者が資料作成・報告書記入等の非商談業務に週30時間以上費やしている"),
        ("顧客情報の分散",   "CRM・メール・議事録・Slackなど複数ツールに情報が散在し、一元管理ができていない"),
        ("予測精度の低さ",   "受注予測は担当者の「勘」に依存。四半期末になるまで着地見通しが不明瞭"),
        ("育成コストの増大", "新人営業の一人前化に平均18ヶ月。ベテランのノウハウが属人化している"),
    ]
    for i, (title, desc) in enumerate(problems):
        y = Inches(1.0 + i * 1.4)
        box = s2.shapes.add_shape(1, Inches(0.5), y, Inches(12), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = PptxColor(255, 240, 240) if i % 2 == 0 else PptxColor(240, 248, 255)
        box.line.color.rgb = PptxColor(200, 200, 200)
        add_text_box(s2, f"課題{i+1}：{title}", Inches(0.7), y+Inches(0.05), Inches(3), Inches(0.45),
                     font_size=14, bold=True, color=(192,0,0))
        add_text_box(s2, desc, Inches(3.8), y+Inches(0.1), Inches(8.8), Inches(0.9),
                     font_size=13, color=(40,40,40))

    # スライド3: 製品概要
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s3, "SalesAI Pro ─ 製品概要", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=22, bold=True, color=(31,56,100))

    features = [
        ("🤖 AI商談アシスト",    "会議録の自動要約・次アクション提案・提案書ドラフト自動生成"),
        ("📊 パイプライン分析",  "機械学習による受注確度スコアリング・受注確度の可視化ダッシュボード"),
        ("🔗 データ統合",        "Salesforce / HubSpot / Gmail / Slack / Zoom との双方向連携"),
        ("📚 ナレッジ共有",      "ベストプラクティスの自動抽出・新人向けレコメンデーション機能"),
    ]
    feat_tbl = s3.shapes.add_table(5, 2, Inches(0.5), Inches(1.0), Inches(12), Inches(4.5)).table
    feat_tbl.cell(0,0).text = "機能"
    feat_tbl.cell(0,1).text = "概要"
    for r in range(2):
        feat_tbl.cell(0,r).fill.solid()
        feat_tbl.cell(0,r).fill.fore_color.rgb = PptxColor(0,70,127)
        feat_tbl.cell(0,r).text_frame.paragraphs[0].runs[0].font.color.rgb = PptxColor(255,255,255)
    for r, (fname, fdesc) in enumerate(features, 1):
        feat_tbl.cell(r,0).text = fname
        feat_tbl.cell(r,1).text = fdesc

    add_text_box(s3, "【想定価格】 月額 180,000円〜（ユーザー数・機能に応じて変動）\n"
                     "【リリース目標】 2024年10月（β版） → 2025年1月（正式リリース）",
                 Inches(0.5), Inches(5.8), Inches(12), Inches(1.4),
                 font_size=13, color=(60,60,60))

    # スライド4: ロードマップ
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s4, "開発ロードマップ", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=22, bold=True, color=(31,56,100))

    road_data = [
        ["フェーズ",   "Phase 1（Q3 2024）",     "Phase 2（Q4 2024）",      "Phase 3（Q1 2025）",     "Phase 4（Q2 2025）"],
        ["テーマ",     "基盤構築・MVP",            "β版リリース",             "正式リリース",            "拡張・グローバル"],
        ["主要機能",   "AI要約・Slack連携",        "スコアリング・ダッシュボ","CRM連携全対応・API公開", "多言語・海外向けUI"],
        ["ターゲット", "社内PoC（10ユーザー）",   "選定パートナー5社",       "国内法人向け一般公開",   "アジア展開開始"],
        ["KPI",        "精度検証・フィードバック", "β満足度85%以上",          "初月100社契約",          "ARR 5億円達成"],
    ]
    road_tbl = s4.shapes.add_table(len(road_data), 5,
                                    Inches(0.3), Inches(1.0), Inches(12.7), Inches(5.5)).table
    for r, row in enumerate(road_data):
        for c, val in enumerate(row):
            road_tbl.cell(r,c).text = val
            if r == 0:
                road_tbl.cell(r,c).fill.solid()
                road_tbl.cell(r,c).fill.fore_color.rgb = PptxColor(31,56,100)
                road_tbl.cell(r,c).text_frame.paragraphs[0].runs[0].font.color.rgb = PptxColor(255,255,255)
            elif c == 0:
                road_tbl.cell(r,c).fill.solid()
                road_tbl.cell(r,c).fill.fore_color.rgb = PptxColor(220,230,242)

    prs.save(DOCS / "新製品企画提案書_SalesAI_Pro.pptx")
    print("  ✓ 新製品企画提案書_SalesAI_Pro.pptx")


# ── PPTX 3: IT中期計画 ──────────────────────
def make_it_plan():
    prs = PRS()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_slide_title(prs, 0, "IT基盤整備 中期計画 2024〜2026",
                    "情報システム部　2024年4月策定")

    # スライド2: 現状課題と目標
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s2, "現状のITシステム課題", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=20, bold=True, color=(31,56,100))
    issues = [
        ("基幹ERPの老朽化",      "2008年導入のオンプレERPがEOLを迎え、2026年以降サポート終了の危機"),
        ("データサイロの深刻化", "部門ごとにExcelと個別システムが乱立。全社横断のデータ分析が不可能"),
        ("セキュリティリスク",   "ゼロトラスト未対応のネットワーク構成。リモートワーク拡大でリスク増大"),
        ("開発・運用の非効率",  "手動デプロイ・テスト。年間リリース回数が競合比で約1/5"),
    ]
    issue_tbl = s2.shapes.add_table(5, 2, Inches(0.5), Inches(0.9), Inches(12.3), Inches(5.5)).table
    issue_tbl.cell(0,0).text = "課題"
    issue_tbl.cell(0,1).text = "内容"
    for c in range(2):
        issue_tbl.cell(0,c).fill.solid()
        issue_tbl.cell(0,c).fill.fore_color.rgb = PptxColor(192,0,0)
        issue_tbl.cell(0,c).text_frame.paragraphs[0].runs[0].font.color.rgb = PptxColor(255,255,255)
    for r, (title, desc) in enumerate(issues, 1):
        issue_tbl.cell(r,0).text = title
        issue_tbl.cell(r,1).text = desc

    # スライド3: 投資計画グラフ
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(s3, "IT投資計画（2024〜2026年度）",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=20, bold=True, color=(31,56,100))

    inv_data = CategoryChartData()
    inv_data.categories = ["2024年度", "2025年度", "2026年度"]
    inv_data.add_series("ERP刷新",        (3.5, 4.2, 1.0))
    inv_data.add_series("クラウド移行",   (1.2, 1.8, 0.8))
    inv_data.add_series("セキュリティ",   (0.8, 1.0, 1.0))
    inv_data.add_series("DX・AI基盤",     (0.5, 1.5, 2.5))

    inv_chart = s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.5), Inches(1.0), Inches(8), Inches(5.5),
        inv_data,
    ).chart
    inv_chart.has_title = True
    inv_chart.chart_title.text_frame.text = "IT投資計画（億円）"

    add_text_box(s3, "投資総額サマリー\n\n"
                     "2024年度：  6.0億円\n"
                     "2025年度：  8.5億円\n"
                     "2026年度：  5.3億円\n"
                     "────────────\n"
                     "3ヶ年合計：19.8億円\n\n"
                     "※ ERP刷新が最大項目。\n  2026年度以降は\n  維持費用のみに移行予定。",
                 Inches(8.8), Inches(1.5), Inches(4.2), Inches(5.0),
                 font_size=14, color=(30,30,30))

    prs.save(DOCS / "IT基盤整備中期計画_2024-2026.pptx")
    print("  ✓ IT基盤整備中期計画_2024-2026.pptx")


# ─────────────────────────────────────────────
#  XLSX
# ─────────────────────────────────────────────
from openpyxl import Workbook
from openpyxl.styles import (Font, PatternFill, Alignment, Border, Side,
                              numbers as xl_num)
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, LineChart, PieChart, Reference, Series
from openpyxl.chart.series import DataPoint


def thick_border():
    s = Side(border_style="thin", color="000000")
    return Border(left=s, right=s, top=s, bottom=s)


def header_style(ws, row, cols, fill_color="1F3864", font_color="FFFFFF"):
    fill = PatternFill("solid", fgColor=fill_color)
    font = Font(bold=True, color=font_color)
    for col in cols:
        cell = ws.cell(row=row, column=col)
        cell.fill = fill
        cell.font = font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = thick_border()


# ── XLSX 1: 月次売上実績表 ──────────────────
def make_sales_xlsx():
    wb = Workbook()

    # ── シート1: 月次サマリ
    ws = wb.active
    ws.title = "月次売上サマリ"

    ws.merge_cells("A1:J1")
    ws["A1"] = "2024年度 月次売上実績表（単位：百万円）"
    ws["A1"].font = Font(bold=True, size=16)
    ws["A1"].alignment = Alignment(horizontal="center")

    ws["A2"] = "作成日：2024年6月30日"
    ws["A2"].font = Font(italic=True, color="888888")

    headers = ["月", "製品事業", "サービス事業", "保守・サポート", "新規事業",
               "売上合計", "予算", "達成率", "前年実績", "前年比"]
    for c, h in enumerate(headers, 1):
        ws.cell(row=4, column=c, value=h)
    header_style(ws, 4, range(1, len(headers)+1))

    months    = ["4月", "5月", "6月", "7月", "8月", "9月", "10月", "11月", "12月", "1月", "2月", "3月"]
    prod_vals = [1820,  1950,  2180,  2050,  2100,  2300,  2150,   2400,   2600,   2200,  2100,  2400]
    svc_vals  = [ 950,   980,  1020,   990,  1010,  1080,  1050,   1120,   1200,   1100,  1050,  1100]
    mnt_vals  = [ 580,   580,   590,   580,   590,   590,   590,    600,    600,    600,   600,   610]
    new_vals  = [ 170,   190,   210,   220,   240,   260,   280,    310,    350,    380,   400,   440]
    budgets   = [5200,  5400,  5700,  5400,  5600,  5900,  5700,   6100,   6400,   5900,  5700,  6200]
    prev_year = [4800,  5100,  5300,  5000,  5200,  5500,  5200,   5700,   5900,   5500,  5300,  5900]

    for i, month in enumerate(months):
        r = i + 5
        total = prod_vals[i] + svc_vals[i] + mnt_vals[i] + new_vals[i]
        ws.cell(r, 1, month)
        ws.cell(r, 2, prod_vals[i])
        ws.cell(r, 3, svc_vals[i])
        ws.cell(r, 4, mnt_vals[i])
        ws.cell(r, 5, new_vals[i])
        ws.cell(r, 6, total)
        ws.cell(r, 7, budgets[i])
        ws.cell(r, 8, round(total / budgets[i], 4))
        ws.cell(r, 8).number_format = "0.0%"
        ws.cell(r, 9, prev_year[i])
        ws.cell(r, 10, round(total / prev_year[i], 4))
        ws.cell(r, 10).number_format = "0.0%"
        if i % 2 == 0:
            for c in range(1, 11):
                ws.cell(r, c).fill = PatternFill("solid", fgColor="EBF3FB")
        for c in range(1, 11):
            ws.cell(r, c).border = thick_border()

    # 合計行
    total_row = 17
    ws.cell(total_row, 1, "合計").font = Font(bold=True)
    for c in range(2, 7):
        ws.cell(total_row, c, f"=SUM({get_column_letter(c)}5:{get_column_letter(c)}16)")
        ws.cell(total_row, c).font = Font(bold=True)
    for c in range(1, 11):
        ws.cell(total_row, c).fill = PatternFill("solid", fgColor="FFD700")
        ws.cell(total_row, c).border = thick_border()

    ws.column_dimensions["A"].width = 8
    for c in range(2, 11):
        ws.column_dimensions[get_column_letter(c)].width = 15

    # 棒グラフ
    bar_chart = BarChart()
    bar_chart.type = "col"
    bar_chart.title = "月次売上高推移（百万円）"
    bar_chart.y_axis.title = "売上高（百万円）"
    bar_chart.x_axis.title = "月"
    bar_chart.grouping = "stacked"
    bar_chart.overlap = 100

    cats = Reference(ws, min_col=1, min_row=5, max_row=16)
    for col, label in [(2,"製品事業"), (3,"サービス事業"), (4,"保守・サポート"), (5,"新規事業")]:
        data = Reference(ws, min_col=col, min_row=4, max_row=16)
        series = Series(data, title_from_data=True)
        bar_chart.append(series)
    bar_chart.set_categories(cats)
    bar_chart.width  = 24
    bar_chart.height = 14
    ws.add_chart(bar_chart, "A19")

    # 折れ線グラフ（達成率）
    line_chart = LineChart()
    line_chart.title = "月次達成率推移"
    line_chart.y_axis.title = "達成率"
    line_chart.x_axis.title = "月"
    line_data = Reference(ws, min_col=8, min_row=4, max_row=16)
    line_chart.add_data(line_data, titles_from_data=True)
    line_chart.set_categories(cats)
    line_chart.width  = 18
    line_chart.height = 10
    ws.add_chart(line_chart, "L19")

    # ── シート2: 得意先別売上
    ws2 = wb.create_sheet("得意先別売上")
    ws2.merge_cells("A1:F1")
    ws2["A1"] = "2024年度上半期（4〜9月）得意先別売上ランキング"
    ws2["A1"].font = Font(bold=True, size=14)
    ws2["A1"].alignment = Alignment(horizontal="center")

    headers2 = ["順位", "得意先名", "売上高（百万円）", "構成比", "前年比", "主要商品"]
    for c, h in enumerate(headers2, 1):
        ws2.cell(row=3, column=c, value=h)
    header_style(ws2, 3, range(1, 7), fill_color="375623")

    clients = [
        (1,  "株式会社グローバルテック",     2840, "11.4%", "+18.2%", "ProSeries X / サポートパック"),
        (2,  "ナショナル産業株式会社",        2150, "8.7%",  "+5.3%",  "製造ラインシステム"),
        (3,  "東日本情報サービス株式会社",    1980, "8.0%",  "+22.1%", "AIサービスパッケージ"),
        (4,  "株式会社マーケットリーダー",    1750, "7.1%",  "-2.4%",  "従来製品（更新待ち）"),
        (5,  "西部製造工業株式会社",          1620, "6.5%",  "+9.8%",  "設備管理ソフトウェア"),
        (6,  "有限会社デジタルソリューション", 980, "4.0%",  "+35.6%", "クラウドサービス"),
        (7,  "株式会社フューチャーネット",     870, "3.5%",  "+11.2%", "ネットワーク機器"),
        (8,  "中央コンサルティング株式会社",  760, "3.1%",  "-8.7%",  "コンサルサービス"),
        (9,  "アジア太平洋ロジスティクス",     640, "2.6%",  "+44.9%", "物流管理システム"),
        (10, "北海道電力関連会社",             590, "2.4%",  "+3.1%",  "エネルギー管理"),
    ]
    for r_offset, row in enumerate(clients):
        r = r_offset + 4
        for c, val in enumerate(row, 1):
            ws2.cell(r, c, val)
            ws2.cell(r, c).border = thick_border()
        if r_offset % 2 == 0:
            for c in range(1, 7):
                ws2.cell(r, c).fill = PatternFill("solid", fgColor="EBF7EE")

    for c, w in enumerate([8, 30, 20, 12, 12, 30], 1):
        ws2.column_dimensions[get_column_letter(c)].width = w

    # 円グラフ
    pie = PieChart()
    pie.title = "得意先別売上構成（上位10社）"
    pie_data = Reference(ws2, min_col=3, min_row=3, max_row=13)
    pie_cats = Reference(ws2, min_col=2, min_row=4, max_row=13)
    pie.add_data(pie_data, titles_from_data=True)
    pie.set_categories(pie_cats)
    pie.width  = 18
    pie.height = 14
    ws2.add_chart(pie, "A16")

    wb.save(DOCS / "月次売上実績表_2024年度.xlsx")
    print("  ✓ 月次売上実績表_2024年度.xlsx")


# ── XLSX 2: プロジェクト工程管理表 ───────────
def make_project_xlsx():
    wb = Workbook()
    ws = wb.active
    ws.title = "工程管理表"

    ws.merge_cells("A1:T1")
    ws["A1"] = "プロジェクト工程管理表　─ 基幹ERPシステム刷新プロジェクト"
    ws["A1"].font = Font(bold=True, size=15, color="1F3864")
    ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 30

    ws.merge_cells("A2:T2")
    ws["A2"] = "プロジェクト期間：2024年4月〜2026年3月　　PM：佐藤 太郎　　ステータス：進行中"
    ws["A2"].font = Font(size=11, italic=True, color="555555")

    # ヘッダー行
    fixed_headers = ["WBS", "フェーズ", "タスク名", "担当部署", "担当者", "計画開始", "計画終了", "進捗%", "ステータス"]
    months_header = ["2024\n4月", "5月", "6月", "7月", "8月", "9月", "10月", "11月", "12月",
                     "2025\n1月", "2月"]

    for c, h in enumerate(fixed_headers, 1):
        cell = ws.cell(row=4, column=c, value=h)
        cell.fill = PatternFill("solid", fgColor="1F3864")
        cell.font = Font(bold=True, color="FFFFFF", size=10)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = thick_border()

    for c, h in enumerate(months_header, len(fixed_headers)+1):
        cell = ws.cell(row=4, column=c, value=h)
        cell.fill = PatternFill("solid", fgColor="375623")
        cell.font = Font(bold=True, color="FFFFFF", size=9)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = thick_border()

    tasks = [
        # WBS, フェーズ, タスク名, 担当, 担当者, 開始, 終了, 進捗, Status, ガント(11列)
        ("1",   "計画",   "現状業務分析・要件定義",     "全部門",   "鈴木PM補佐", "4/1",  "5/31",  100, "完了",    [1,1,0,0,0,0,0,0,0,0,0]),
        ("1.1", "計画",   "AS-ISプロセスマッピング",    "業務改革室","山田 花",    "4/1",  "4/30",  100, "完了",    [1,0,0,0,0,0,0,0,0,0,0]),
        ("1.2", "計画",   "TO-BE要件定義書作成",        "業務改革室","田中 一",    "5/1",  "5/31",  100, "完了",    [0,1,0,0,0,0,0,0,0,0,0]),
        ("2",   "選定",   "ベンダー選定・RFP対応",      "購買/IT",  "高橋 二",    "5/15", "7/31",  85,  "進行中",  [0,1,1,1,0,0,0,0,0,0,0]),
        ("2.1", "選定",   "RFP作成・送付（5社）",       "IT部門",   "木村 三",    "5/15", "6/15",  100, "完了",    [0,1,1,0,0,0,0,0,0,0,0]),
        ("2.2", "選定",   "ベンダープレゼン・評価",      "評価委員会","全委員",     "6/20", "7/20",  60,  "進行中",  [0,0,1,1,0,0,0,0,0,0,0]),
        ("2.3", "選定",   "最終選定・契約締結",          "役員会議", "社長承認",   "7/21", "7/31",  0,   "未着手",  [0,0,0,1,0,0,0,0,0,0,0]),
        ("3",   "設計",   "システム設計・カスタマイズ定義","IT/業務",  "佐藤 四",    "8/1",  "11/30", 0,   "未着手",  [0,0,0,0,1,1,1,1,0,0,0]),
        ("3.1", "設計",   "基本設計書作成",              "ITベンダー","PJ担当",     "8/1",  "9/30",  0,   "未着手",  [0,0,0,0,1,1,0,0,0,0,0]),
        ("3.2", "設計",   "インターフェース設計",         "IT部門",   "木村 五",    "10/1", "10/31", 0,   "未着手",  [0,0,0,0,0,0,1,0,0,0,0]),
        ("3.3", "設計",   "データ移行設計",               "IT部門",   "高橋 六",    "10/1", "11/30", 0,   "未着手",  [0,0,0,0,0,0,1,1,0,0,0]),
        ("4",   "開発",   "システム開発・単体テスト",    "ITベンダー","開発チーム",  "10/1", "2/28",  0,   "未着手",  [0,0,0,0,0,0,1,1,1,1,1]),
        ("5",   "テスト", "統合テスト・ユーザー受入",    "全部門",   "各部門長",   "1/6",  "2/28",  0,   "未着手",  [0,0,0,0,0,0,0,0,0,1,1]),
    ]

    status_colors = {"完了": "C6EFCE", "進行中": "FFEB9C", "未着手": "F2F2F2", "遅延": "FFC7CE"}
    gantt_color   = "4472C4"

    for r_offset, task in enumerate(tasks):
        r = r_offset + 5
        wbs, phase, name, dept, person, start, end, pct, status, gantt = task
        row_data = [wbs, phase, name, dept, person, start, end, f"{pct}%", status]
        for c, val in enumerate(row_data, 1):
            cell = ws.cell(r, c, val)
            cell.border = thick_border()
            cell.alignment = Alignment(vertical="center", wrap_text=True)
            if c == 9:  # ステータス
                cell.fill = PatternFill("solid", fgColor=status_colors.get(status, "F2F2F2"))
                cell.alignment = Alignment(horizontal="center", vertical="center")
            if r_offset % 2 == 0 and c != 9:
                cell.fill = PatternFill("solid", fgColor="F5F5F5")
        for g, flag in enumerate(gantt):
            c = len(fixed_headers) + 1 + g
            cell = ws.cell(r, c)
            cell.border = thick_border()
            if flag:
                cell.fill = PatternFill("solid", fgColor=gantt_color if status != "完了" else "70AD47")
                cell.value = "▬"
                cell.alignment = Alignment(horizontal="center", vertical="center")

    # 列幅
    col_widths = [6, 10, 28, 14, 12, 8, 8, 8, 10] + [5]*11
    for c, w in enumerate(col_widths, 1):
        ws.column_dimensions[get_column_letter(c)].width = w
    ws.row_dimensions[4].height = 35
    ws.freeze_panes = "J5"

    # ── シート2: リスク管理表
    ws2 = wb.create_sheet("リスク管理")
    ws2.merge_cells("A1:G1")
    ws2["A1"] = "リスク管理表"
    ws2["A1"].font = Font(bold=True, size=14)
    ws2["A1"].alignment = Alignment(horizontal="center")

    risk_headers = ["No.", "リスク内容", "発生確率", "影響度", "リスクスコア", "対応策", "担当"]
    for c, h in enumerate(risk_headers, 1):
        cell = ws2.cell(2, c, h)
        cell.fill = PatternFill("solid", fgColor="C00000")
        cell.font = Font(bold=True, color="FFFFFF")
        cell.alignment = Alignment(horizontal="center")
        cell.border = thick_border()

    risks = [
        (1, "ベンダー選定の遅延",           "中",  "高",  15, "選定基準の事前明確化・代替ベンダーリスト準備", "高橋二"),
        (2, "データ移行における品質問題",    "中",  "高",  15, "移行テストを3段階実施、データクレンジング前倒し","木村三"),
        (3, "ユーザー教育の遅れ",           "高",  "中",  12, "e-ラーニング導入・部門ごとのトレーナー任命",    "田中一"),
        (4, "予算超過",                     "低",  "高",  9,  "月次予算レビュー・変更管理プロセスの厳格化",    "鈴木補佐"),
        (5, "カスタマイズの過多",           "高",  "中",  12, "Fit/Gap分析でカスタマイズを最小化方針徹底",     "佐藤四"),
        (6, "キーパーソンの離脱",           "低",  "高",  9,  "バックアップ要員の確保・ナレッジ文書化",        "PM全体"),
    ]
    risk_fill = {"高": "FFC7CE", "中": "FFEB9C", "低": "C6EFCE"}
    for r_offset, risk in enumerate(risks):
        r = r_offset + 3
        for c, val in enumerate(risk, 1):
            cell = ws2.cell(r, c, val)
            cell.border = thick_border()
            cell.alignment = Alignment(vertical="center", wrap_text=True)
            if c in (3, 4):
                cell.fill = PatternFill("solid", fgColor=risk_fill.get(str(val), "FFFFFF"))
                cell.alignment = Alignment(horizontal="center", vertical="center")

    for c, w in enumerate([5, 30, 12, 10, 14, 40, 12], 1):
        ws2.column_dimensions[get_column_letter(c)].width = w

    wb.save(DOCS / "プロジェクト工程管理表_ERP刷新PJ.xlsx")
    print("  ✓ プロジェクト工程管理表_ERP刷新PJ.xlsx")


# ─────────────────────────────────────────────
#  PDF (reportlab)
# ─────────────────────────────────────────────
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm, mm
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                  TableStyle, HRFlowable, PageBreak)
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT


def setup_pdf_fonts():
    pdfmetrics.registerFont(UnicodeCIDFont("HeiseiKakuGo-W5"))
    pdfmetrics.registerFont(UnicodeCIDFont("HeiseiMin-W3"))


def make_styles():
    styles = getSampleStyleSheet()
    base = "HeiseiKakuGo-W5"
    custom = {
        "title":    ParagraphStyle("title",    fontName=base, fontSize=20, alignment=TA_CENTER, spaceAfter=6),
        "subtitle": ParagraphStyle("subtitle", fontName=base, fontSize=12, alignment=TA_CENTER, textColor=colors.HexColor("#555555"), spaceAfter=16),
        "h1":       ParagraphStyle("h1",       fontName=base, fontSize=14, textColor=colors.HexColor("#1F3864"), spaceBefore=14, spaceAfter=6),
        "h2":       ParagraphStyle("h2",       fontName=base, fontSize=12, textColor=colors.HexColor("#2E74B5"), spaceBefore=10, spaceAfter=4),
        "body":     ParagraphStyle("body",     fontName="HeiseiMin-W3", fontSize=10, leading=16, spaceAfter=6),
        "note":     ParagraphStyle("note",     fontName=base, fontSize=9,  textColor=colors.HexColor("#666666"), spaceAfter=4),
        "center":   ParagraphStyle("center",   fontName=base, fontSize=10, alignment=TA_CENTER),
    }
    return custom


def tbl_style_base(header_color="#1F3864"):
    hc = colors.HexColor(header_color)
    return TableStyle([
        ("BACKGROUND",  (0,0), (-1,0), hc),
        ("TEXTCOLOR",   (0,0), (-1,0), colors.white),
        ("FONTNAME",    (0,0), (-1,0), "HeiseiKakuGo-W5"),
        ("FONTSIZE",    (0,0), (-1,-1), 9),
        ("FONTNAME",    (0,1), (-1,-1), "HeiseiMin-W3"),
        ("ALIGN",       (0,0), (-1,-1), "CENTER"),
        ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
        ("GRID",        (0,0), (-1,-1), 0.5, colors.grey),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#EBF3FB")]),
        ("TOPPADDING",  (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ])


# ── PDF 1: 会社概要 ─────────────────────────
def make_company_profile_pdf():
    setup_pdf_fonts()
    s = make_styles()
    path = str(DOCS / "会社概要_株式会社テックリードジャパン.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4,
                            rightMargin=2*cm, leftMargin=2*cm,
                            topMargin=2.5*cm, bottomMargin=2*cm)
    story = []

    story.append(Spacer(1, 1*cm))
    story.append(Paragraph("会 社 概 要", s["title"]))
    story.append(Paragraph("株式会社テックリードジャパン　Company Profile 2024", s["subtitle"]))
    story.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor("#1F3864")))
    story.append(Spacer(1, 0.5*cm))

    story.append(Paragraph("■ 基本情報", s["h1"]))
    basic_data = [
        ["項目", "内容"],
        ["商号",       "株式会社テックリードジャパン"],
        ["英文名",     "TechLead Japan Co., Ltd."],
        ["設立",       "2005年（平成17年）4月1日"],
        ["資本金",     "5億円"],
        ["代表者",     "代表取締役社長　田中 一郎"],
        ["従業員数",   "850名（2024年4月1日現在、連結）"],
        ["本社所在地", "〒100-0001　東京都千代田区丸の内2-1-1 テックビル15F"],
        ["売上高",     "140億円（2024年3月期　連結）"],
        ["事業内容",   "ITソリューション事業・AIサービス事業・保守サポート事業"],
        ["グループ会社","株式会社テックラボ（100%子会社）　他2社"],
        ["取引銀行",   "みずほ銀行　丸の内支店 / 三菱UFJ銀行 大手町支店"],
        ["主要顧客",   "製造業・流通業・金融機関　他（上場企業200社超）"],
        ["認証取得",   "ISO 9001:2015 / ISO 27001:2013 / プライバシーマーク"],
    ]
    t = Table(basic_data, colWidths=[4.5*cm, 11.5*cm])
    t.setStyle(tbl_style_base())
    t.setStyle(TableStyle([("ALIGN", (0,0), (0,-1), "CENTER"),
                            ("ALIGN", (1,1), (1,-1), "LEFT"),
                            ("LEFTPADDING", (1,1), (1,-1), 8)]))
    story.append(t)

    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph("■ 事業部門別売上構成（2024年3月期）", s["h1"]))
    seg_data = [
        ["事業部門", "売上高", "構成比", "前年比", "主要サービス"],
        ["ITソリューション事業", "84.0億円", "60%", "+10.5%", "ERP・CRM導入支援、システム開発"],
        ["AIサービス事業",       "28.0億円", "20%", "+35.2%", "AI・機械学習基盤、データ分析"],
        ["保守サポート事業",     "21.0億円", "15%",  "+3.1%", "システム運用保守、ヘルプデスク"],
        ["新規事業",             "7.0億円",  "5%",  "+165%",  "SaaS型営業支援ツール"],
        ["合計",                "140.0億円", "100%", "+12.0%", "─"],
    ]
    t2 = Table(seg_data, colWidths=[3.8*cm, 2.5*cm, 1.8*cm, 1.8*cm, 6.1*cm])
    t2.setStyle(tbl_style_base())
    story.append(t2)

    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph("■ 沿革", s["h1"]))
    history = [
        ["年月",       "主な出来事"],
        ["2005年4月", "東京都千代田区にて資本金3,000万円で創業。ERP導入支援事業を開始"],
        ["2008年9月", "本社を現在地（丸の内）に移転。従業員数100名を突破"],
        ["2011年3月", "東日本大震災BCP対応として大阪・福岡にサブオフィスを設置"],
        ["2014年7月", "ISO 27001（情報セキュリティ）認証取得"],
        ["2017年4月", "AI研究部門を設立。機械学習・自然言語処理の事業化に着手"],
        ["2019年10月", "AIサービス部門を独立させ、子会社「テックラボ」を設立"],
        ["2021年4月", "プライバシーマーク取得。個人情報保護体制を強化"],
        ["2022年11月", "SaaS型営業支援ツール「SalesAI Pro」β版リリース"],
        ["2024年4月", "資本金を5億円に増資。従業員850名体制へ拡大"],
    ]
    t3 = Table(history, colWidths=[3*cm, 13*cm])
    t3.setStyle(tbl_style_base())
    t3.setStyle(TableStyle([("ALIGN", (1,1), (1,-1), "LEFT"),
                             ("LEFTPADDING", (1,1), (1,-1), 8)]))
    story.append(t3)

    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph("■ 組織図（概要）", s["h1"]))
    org_data = [
        ["代表取締役社長", "", "", "", ""],
        ["取締役会・監査役会", "", "経営企画室", "内部監査室", ""],
        ["営業本部", "技術本部", "AI事業本部", "管理本部", "新規事業本部"],
        ["国内営業部\nグローバル営業部\nインサイドセールス",
         "システム開発部\nインフラ・SRE部\nQA部",
         "AIエンジニアリング部\nデータサイエンス部\nプロダクト部",
         "経理財務部\n人事部\n総務・法務部",
         "SaaS事業部\n海外事業準備室"],
    ]
    t4 = Table(org_data, colWidths=[3.2*cm]*5)
    t4.setStyle(TableStyle([
        ("SPAN",      (0,0), (4,0)),
        ("BACKGROUND",(0,0),(4,0), colors.HexColor("#1F3864")),
        ("TEXTCOLOR", (0,0),(4,0), colors.white),
        ("ALIGN",     (0,0),(-1,-1), "CENTER"),
        ("VALIGN",    (0,0),(-1,-1), "MIDDLE"),
        ("FONTNAME",  (0,0),(-1,-1), "HeiseiKakuGo-W5"),
        ("FONTSIZE",  (0,0),(-1,-1), 9),
        ("GRID",      (0,0),(-1,-1), 0.5, colors.grey),
        ("TOPPADDING",(0,0),(-1,-1), 6),
        ("BOTTOMPADDING",(0,0),(-1,-1), 6),
        ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.HexColor("#D6E4F0"), colors.HexColor("#EBF3FB"), colors.white]),
    ]))
    story.append(t4)

    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph("■ 主要な受賞・認定", s["h1"]))
    awards = [
        ["年",    "受賞・認定内容",                                        "授与機関"],
        ["2019", "DX認定事業者",                                           "経済産業省"],
        ["2021", "健康経営優良法人（大規模法人部門）",                     "経済産業省・日本健康会議"],
        ["2022", "ITreview Grid Award 2022 Winter（SaaS営業支援カテゴリ）", "ITreview"],
        ["2023", "女性活躍推進企業認定（えるぼし三つ星）",                  "厚生労働省"],
        ["2024", "中小企業庁 IT活用先進企業100選",                         "中小企業庁"],
    ]
    t5 = Table(awards, colWidths=[1.5*cm, 9.5*cm, 5*cm])
    t5.setStyle(tbl_style_base())
    story.append(t5)

    doc.build(story)
    print("  ✓ 会社概要_株式会社テックリードジャパン.pdf")


# ── PDF 2: 顧客満足度調査報告書 ─────────────
def make_cs_survey_pdf():
    setup_pdf_fonts()
    s = make_styles()
    path = str(DOCS / "顧客満足度調査報告書_2024年春.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4,
                            rightMargin=2*cm, leftMargin=2*cm,
                            topMargin=2.5*cm, bottomMargin=2*cm)
    story = []

    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph("2024年春季　顧客満足度調査報告書", s["title"]))
    story.append(Paragraph("調査実施期間：2024年4月1日〜4月30日　　回答数：487社（回答率68.3%）", s["subtitle"]))
    story.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor("#1F3864")))
    story.append(Spacer(1, 0.3*cm))

    story.append(Paragraph("1. 調査概要", s["h1"]))
    overview_data = [
        ["調査目的", "製品・サービスの品質向上および顧客リテンション向上のための実態把握"],
        ["調査対象", "直近12ヶ月以内に取引のある法人顧客（国内）"],
        ["調査方法", "Webアンケート（SurveyMonkey）+ 一部電話ヒアリング"],
        ["調査期間", "2024年4月1日（月）〜4月30日（火）"],
        ["送付数",   "713社　有効回答：487社（回答率68.3%）"],
        ["集計日",   "2024年5月10日"],
        ["調査担当", "マーケティング部　CS推進チーム"],
    ]
    t = Table(overview_data, colWidths=[3.5*cm, 12.5*cm])
    t.setStyle(tbl_style_base("#375623"))
    t.setStyle(TableStyle([("ALIGN",(1,1),(1,-1),"LEFT"), ("LEFTPADDING",(1,1),(1,-1),8)]))
    story.append(t)

    story.append(Spacer(1, 0.4*cm))
    story.append(Paragraph("2. 総合満足度スコア（NPS・CSATサマリ）", s["h1"]))
    score_data = [
        ["指標", "今回（2024春）", "前回（2023秋）", "前々回（2023春）", "目標値", "業界平均"],
        ["CSAT（総合満足度）",   "82.4%",  "79.8%",  "76.5%",  "80.0%", "74.2%"],
        ["NPS（推奨意向）",      "+28",    "+22",    "+18",    "+25",   "+15"  ],
        ["製品品質満足度",        "85.1%",  "82.3%",  "79.1%",  "83.0%", "76.8%"],
        ["サポート満足度",        "78.9%",  "75.2%",  "70.8%",  "78.0%", "68.5%"],
        ["価格妥当性",            "64.3%",  "65.1%",  "63.2%",  "65.0%", "61.0%"],
        ["継続利用意向",          "89.2%",  "86.7%",  "84.1%",  "88.0%", "79.3%"],
    ]
    t2 = Table(score_data, colWidths=[3.8*cm, 2.2*cm, 2.2*cm, 2.2*cm, 2.0*cm, 2.0*cm])
    t2.setStyle(tbl_style_base())
    story.append(t2)

    story.append(Spacer(1, 0.4*cm))
    story.append(Paragraph("3. 設問別スコア詳細", s["h1"]))

    story.append(Paragraph("3-1. 製品・サービス品質", s["h2"]))
    q_data = [
        ["設問", "5点\n（大変満足）", "4点\n（満足）", "3点\n（普通）", "2点\n（やや不満）", "1点\n（不満）", "平均スコア"],
        ["Q1. 製品の機能充実度",        "38.4%", "46.7%", "11.5%",  "2.7%",  "0.7%", "4.19"],
        ["Q2. 製品の安定性・信頼性",    "42.1%", "43.9%", "10.5%",  "2.9%",  "0.6%", "4.24"],
        ["Q3. 操作のしやすさ（UI/UX）", "29.6%", "44.2%", "18.3%",  "6.4%",  "1.5%", "3.94"],
        ["Q4. 導入・セットアップの容易さ","25.1%","43.7%", "20.5%",  "8.4%",  "2.3%", "3.81"],
        ["Q5. アップデート・機能追加",  "31.8%", "42.5%", "19.1%",  "5.5%",  "1.1%", "3.98"],
    ]
    t3 = Table(q_data, colWidths=[4.5*cm, 1.7*cm, 1.7*cm, 1.7*cm, 1.7*cm, 1.7*cm, 2.0*cm])
    t3.setStyle(tbl_style_base())
    story.append(t3)

    story.append(Paragraph("3-2. サポート・対応品質", s["h2"]))
    q2_data = [
        ["設問", "5点", "4点", "3点", "2点", "1点", "平均"],
        ["Q6. 問い合わせへの応答速度",   "28.5%", "42.7%", "19.2%",  "7.4%",  "2.2%", "3.88"],
        ["Q7. 問題解決までの適切さ",      "31.4%", "45.8%", "15.6%",  "5.7%",  "1.5%", "4.00"],
        ["Q8. 担当者の技術知識",          "35.9%", "44.6%", "14.4%",  "4.1%",  "1.0%", "4.10"],
        ["Q9. 担当者のコミュニケーション", "38.2%", "43.9%", "13.7%",  "3.3%",  "0.9%", "4.16"],
        ["Q10. ドキュメント・マニュアルの充実度","22.3%","38.7%","25.8%","10.1%","3.1%","3.66"],
    ]
    t4 = Table(q2_data, colWidths=[4.5*cm, 1.7*cm, 1.7*cm, 1.7*cm, 1.7*cm, 1.7*cm, 2.0*cm])
    t4.setStyle(tbl_style_base("#375623"))
    story.append(t4)

    story.append(PageBreak())
    story.append(Paragraph("4. 自由回答コメント（抜粋）", s["h1"]))
    story.append(Paragraph("4-1. 高評価コメント（ポジティブ）", s["h2"]))
    pos_comments = [
        ["回答者属性", "コメント"],
        ["製造業・大手（500名以上）",
         "「AI機能の精度が想定以上に高く、提案書自動生成機能で営業の工数が週5時間削減できた。」"],
        ["金融業・中堅（50〜500名）",
         "「サポートチームのレスポンスが常に早く、問題が翌日以内に解決される点に満足している。」"],
        ["流通業・中小（50名以下）",
         "「直感的なUIで、IT部門なしの自社導入が実現できた。費用対効果も高い。」"],
        ["サービス業・大手",
         "「定期的なバージョンアップで常に最新機能が使えるのが魅力。ロードマップの透明性も良い。」"],
    ]
    t5 = Table(pos_comments, colWidths=[4.5*cm, 11.5*cm])
    t5.setStyle(tbl_style_base())
    t5.setStyle(TableStyle([("ALIGN",(1,1),(1,-1),"LEFT"), ("LEFTPADDING",(1,1),(1,-1),6)]))
    story.append(t5)

    story.append(Spacer(1, 0.3*cm))
    story.append(Paragraph("4-2. 改善要望コメント（要対応）", s["h2"]))
    neg_comments = [
        ["回答者属性", "コメント", "重要度"],
        ["製造業・中堅",  "「モバイルアプリの操作性がWeb版と乖離しており、現場利用に支障がある。」",      "高"],
        ["流通業・大手",  "「価格が競合と比べて高めに感じる。ボリュームディスカウントの拡充を要望する。」", "中"],
        ["IT業・中小",    "「APIドキュメントの更新がリリースに追いついていないことが多い。」",              "高"],
        ["金融業・中堅",  "「障害発生時の情報提供が遅い。ステータスページの充実を希望する。」",            "高"],
        ["医療・中小",    "「医療業界特有のコンプライアンス要件対応の機能が不足している。」",              "中"],
    ]
    t6 = Table(neg_comments, colWidths=[3.5*cm, 10.0*cm, 2.5*cm])
    t6.setStyle(tbl_style_base("#C00000"))
    t6.setStyle(TableStyle([("ALIGN",(1,1),(1,-1),"LEFT"), ("LEFTPADDING",(1,1),(1,-1),6)]))
    story.append(t6)

    story.append(Spacer(1, 0.4*cm))
    story.append(Paragraph("5. アクションプラン", s["h1"]))
    actions = [
        ["優先度", "課題",              "対応策",                       "担当部門",    "完了目標"],
        ["緊急",  "モバイルアプリUX",  "UI/UXリニューアル（v3.2）",    "プロダクト部","2024年9月"],
        ["高",    "APIドキュメント遅延","ドキュメントチーム専任1名増員","技術部",      "2024年7月"],
        ["高",    "障害情報提供の遅れ","ステータスページ自動更新化",    "SRE部",       "2024年8月"],
        ["中",    "価格競争力",        "ボリュームディスカウント改定",   "営業企画部",  "2024年10月"],
        ["中",    "業界特化機能",      "医療・金融向けオプション検討",   "PdM",         "2025年Q1"],
    ]
    t7 = Table(actions, colWidths=[1.5*cm, 3.5*cm, 5.0*cm, 3.0*cm, 2.5*cm])
    t7.setStyle(tbl_style_base())
    story.append(t7)

    doc.build(story)
    print("  ✓ 顧客満足度調査報告書_2024年春.pdf")


# ─────────────────────────────────────────────
if __name__ == "__main__":
    print("=== ダミードキュメント生成開始 ===\n")

    print("[DOCX]")
    make_board_minutes()
    make_manual()
    make_recruitment()

    print("\n[PPTX]")
    make_q1_report()
    make_product_proposal()
    make_it_plan()

    print("\n[XLSX]")
    make_sales_xlsx()
    make_project_xlsx()

    print("\n[PDF]")
    make_company_profile_pdf()
    make_cs_survey_pdf()

    print(f"\n=== 完了: docs/ に {len(list(Path('docs').iterdir()))} ファイル生成 ===")
